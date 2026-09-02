#!/usr/bin/env python3
"""Перенести слайды из одной .pptx в другую, сохранив фигуры, картинки и заметки.

ЗАЧЕМ. Дека собирается генератором из YAML, но человек правит готовый pptx руками, и когда
ручных линий становится две, слить их нечем: в `preza_merge` backend `graft` намеренно не
реализован (src/preza_merge/apply.py) — он разбирает ФОРМАТИРОВАНИЕ, а не переносит слайды.
Этот скрипт закрывает ровно ту дыру: разовый перенос отдельных слайдов между двумя pptx,
которые выросли из одного шаблона.

ГРАНИЦЫ. Скрипт рассчитан на файлы С ОБЩИМ ПРЕДКОМ: макет ищется по ИМЕНИ, и если имени нет
в целевой деке — падаем, а не подставляем молча другой макет (фигуры уехали бы по месту).
Тема, размер слайда и набор макетов не переносятся вообще.

ИСПОЛЬЗОВАНИЕ
  graft_slides.py <target.pptx> <source.pptx> <out.pptx> [--insert S:AFTER] [--replace S:T] [--report]

    --insert  S:AFTER   вставить слайд S исходника ПОСЛЕ слайда AFTER цели (AFTER=0 — в начало)
    --replace S:T       заменить слайд T цели слайдом S исходника
    --report            показать план и ничего не писать

Номера слайдов везде 1-based и относятся к ИСХОДНОМУ состоянию файлов: порядок операций
внутри скрипта на них не влияет, считать «с поправкой на уже вставленное» не нужно.
"""

import argparse
import copy
import hashlib
import sys

from pptx import Presentation
from pptx.oxml.ns import qn


def _layout_by_name(prs, name):
    """Макет цели с тем же именем. Ошибка, если имени нет — молча подставлять чужой нельзя."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    have = sorted({lay.name for m in prs.slide_masters for lay in m.slide_layouts})
    raise SystemExit(f"graft: макета {name!r} нет в целевой деке; есть: {have}")


def _media_index(part):
    """sha256 → rId уже связанных с частью картинок, чтобы не плодить дубли ppt/media/**."""
    idx = {}
    for rid, rel in part.rels.items():
        if rel.is_external:
            continue
        target = rel.target_part
        if target.partname.startswith("/ppt/media/"):
            idx[hashlib.sha256(target.blob).hexdigest()] = rid
    return idx


def _remap_media(src_slide, dst_slide, stats):
    """Перевесить r:embed/r:link скопированных фигур на части целевой деки.

    Дедуп по содержимому: если такая же картинка в цели уже есть, переиспользуем её часть.
    """
    dst_part = dst_slide.part
    known = _media_index(dst_part)
    for attr in (qn("r:embed"), qn("r:link")):
        for el in dst_slide.shapes._spTree.iter():
            rid = el.get(attr)
            if not rid:
                continue
            try:
                src_media = src_slide.part.related_part(rid)
            except KeyError:
                raise SystemExit(f"graft: в исходном слайде нет связи {rid} — файл повреждён")
            digest = hashlib.sha256(src_media.blob).hexdigest()
            if digest in known:
                stats["reused"] += 1
            else:
                known[digest] = dst_part.relate_to(src_media, src_media.rels[rid].reltype
                                                   if rid in src_media.rels else
                                                   src_slide.part.rels[rid].reltype)
                stats["added"] += 1
            el.set(attr, known[digest])


def _copy_slide(prs, src_slide):
    """Добавить в конец prs копию src_slide и вернуть её."""
    layout = _layout_by_name(prs, src_slide.slide_layout.name)
    new = prs.slides.add_slide(layout)

    # add_slide копирует плейсхолдеры макета — они не нужны, содержимое приходит целиком.
    spTree = new.shapes._spTree
    for shape in list(spTree):
        if shape.tag in (qn("p:sp"), qn("p:pic"), qn("p:graphicFrame"),
                         qn("p:grpSp"), qn("p:cxnSp")):
            spTree.remove(shape)
    for shape in src_slide.shapes._spTree:
        if shape.tag in (qn("p:sp"), qn("p:pic"), qn("p:graphicFrame"),
                         qn("p:grpSp"), qn("p:cxnSp")):
            spTree.append(copy.deepcopy(shape))

    if src_slide.has_notes_slide:
        src_tf = src_slide.notes_slide.notes_text_frame
        dst_notes = new.notes_slide
        dst_body = dst_notes.notes_text_frame._txBody
        for para in list(dst_body.findall(qn("a:p"))):
            dst_body.remove(para)
        for para in src_tf._txBody.findall(qn("a:p")):
            dst_body.append(copy.deepcopy(para))
    return new


def _move(prs, sld_id_el, index):
    """Переставить уже добавленный sldId на позицию index (0-based) в sldIdLst."""
    lst = prs.slides._sldIdLst
    lst.remove(sld_id_el)
    lst.insert(index, sld_id_el)


def _drop(prs, index):
    """Убрать слайд по 0-based индексу вместе с его частью."""
    lst = prs.slides._sldIdLst
    sld_id = lst[index]
    prs.part.drop_rel(sld_id.rId)
    lst.remove(sld_id)


def _pair(text, flag):
    try:
        left, right = text.split(":")
        return int(left), int(right)
    except ValueError:
        raise SystemExit(f"graft: {flag} ждёт вид N:M, получено {text!r}")


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("target")
    ap.add_argument("source")
    ap.add_argument("out")
    ap.add_argument("--insert", action="append", default=[], metavar="SRC:AFTER")
    ap.add_argument("--replace", action="append", default=[], metavar="SRC:TGT")
    ap.add_argument("--report", action="store_true", help="показать план и выйти")
    args = ap.parse_args(argv)

    tgt = Presentation(args.target)
    src = Presentation(args.source)
    src_slides = list(src.slides)
    n_tgt = len(tgt.slides._sldIdLst)

    def title(slide):
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                return sh.text_frame.text.strip().splitlines()[0][:64]
        return "(без текста)"

    inserts = [_pair(x, "--insert") for x in args.insert]
    replaces = [_pair(x, "--replace") for x in args.replace]
    for s, _ in inserts + replaces:
        if not 1 <= s <= len(src_slides):
            raise SystemExit(f"graft: в исходнике {len(src_slides)} слайдов, запрошен {s}")
    for _, a in inserts:
        if not 0 <= a <= n_tgt:
            raise SystemExit(f"graft: --insert после слайда {a}, а в цели их {n_tgt}")
    for _, t in replaces:
        if not 1 <= t <= n_tgt:
            raise SystemExit(f"graft: --replace слайда {t}, а в цели их {n_tgt}")

    print(f"target: {args.target}  ({n_tgt} слайдов)")
    print(f"source: {args.source}  ({len(src_slides)} слайдов)")
    for s, a in inserts:
        print(f"  insert  src#{s:>3} → после #{a:<3}  {title(src_slides[s - 1])}")
    for s, t in replaces:
        print(f"  replace src#{s:>3} → вместо #{t:<3}  {title(src_slides[s - 1])}")
    if args.report:
        return 0

    stats = {"added": 0, "reused": 0}
    # Каждой операции — целевая позиция в ИСХОДНОЙ нумерации цели. Сначала копируем все
    # слайды в хвост, потом расставляем: так номера в аргументах не «плывут» по ходу.
    planned = [(a, s, None) for s, a in inserts] + [(t - 1, s, t) for s, t in replaces]
    made = []
    for at, s, drop_at in sorted(planned, key=lambda x: x[0]):
        new = _copy_slide(tgt, src_slides[s - 1])
        _remap_media(src_slides[s - 1], new, stats)
        made.append((at, tgt.slides._sldIdLst[-1], drop_at))

    # Расставляем от конца к началу: правки не сдвигают ещё не обработанные позиции.
    # Ключ включает порядковый номер, чтобы РАЗВЕРНУТЬ и группу с одинаковым `at`: каждая
    # следующая вставка в ту же позицию отодвигает предыдущую, и без разворота несколько
    # слайдов «после #21» легли бы задом наперёд.
    order = sorted(enumerate(made), key=lambda x: (x[1][0], x[0]), reverse=True)
    for _, (at, sld_id, _) in order:
        _move(tgt, sld_id, at)
    # Удаляем заменяемые оригиналы: их индекс сдвинулся ровно на число вставок до него.
    for at, _, drop_at in sorted(made, key=lambda x: x[0], reverse=True):
        if drop_at is not None:
            shift = sum(1 for a2, _, _ in made if a2 <= at)
            _drop(tgt, drop_at - 1 + shift)

    tgt.save(args.out)
    print(f"out:    {args.out}  ({len(Presentation(args.out).slides._sldIdLst)} слайдов)")
    print(f"  медиа: переиспользовано {stats['reused']}, добавлено новых частей {stats['added']}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
