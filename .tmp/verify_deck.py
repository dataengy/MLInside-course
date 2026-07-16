#!/usr/bin/env python3
"""Verify a generated deck: slide/image/notes/materials counts, note-emphasis runs, author check.

Usage: python3 .tmp/verify_deck.py [path.pptx]
"""
import re
import sys

from pptx import Presentation

path = sys.argv[1] if len(sys.argv) > 1 else "data/generated/MLInside_Введение-в-dbt_v3.1.pptx"
p = Presentation(path)
imgs = notes = mats = bold = under = ital = 0
for s in p.slides:
    imgs += any(sh.shape_type == 13 for sh in s.shapes)                       # PICTURE
    mats += any(sh.has_text_frame and "Материалы" in sh.text_frame.text for sh in s.shapes)
    if s.has_notes_slide:
        tf = s.notes_slide.notes_text_frame
        if tf.text.strip():
            notes += 1
        for para in tf.paragraphs:
            for r in para.runs:
                bold += bool(r.font.bold)
                under += bool(r.font.underline)
                ital += bool(r.font.italic)
print(path)
print(f"  slides={len(p.slides._sldIdLst)} images={imgs} notes={notes} materials={mats}")
print(f"  note runs: bold={bold} underline={under} italic={ital}")
blob = open(path, "rb").read().decode("latin-1")
print("  author strings:", bool(re.search(r"Крупий|hnkovr|NikolayKrupiy", blob)))
