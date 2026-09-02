#!/usr/bin/env python3
"""Убрать фигуры со слайда готовой .pptx — по индексу, с обязательным показом того, что уходит.

ЗАЧЕМ. У ручной линии деки контент-YAML нет, и мусор на слайде (остаток прежней вёрстки,
забытая панель поверх новой схемы) убрать нечем, кроме PowerPoint. Руками — значит вне git и
без следа в истории. Скрипт делает то же самое воспроизводимо: команда с индексами остаётся
в Justfile, и через полгода видно, что именно и почему сняли.

ОСТОРОЖНО С ИНДЕКСАМИ. Индекс фигуры — не её свойство, а позиция в списке: убрали одну —
следующие сдвинулись. Поэтому все индексы разбираются от ИСХОДНОГО состояния слайда, а
удаление идёт с конца. И поэтому же скрипт всегда печатает, что удаляет: индекс, взятый из
чужого отчёта или из прошлого прогона, может указывать уже не туда.

ИСПОЛЬЗОВАНИЕ
    drop_shapes.py <deck.pptx> --slide N --shape I [--shape J …] [--out OUT.pptx] [--report]
"""

from __future__ import annotations

import argparse
import sys

from pptx import Presentation

_EMU_IN = 914400


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("deck")
    ap.add_argument("--slide", type=int, required=True, help="номер слайда, 1-based")
    ap.add_argument("--shape", type=int, action="append", required=True,
                    help="индекс фигуры, 0-based, от исходного состояния слайда")
    ap.add_argument("--out")
    ap.add_argument("--report", action="store_true")
    args = ap.parse_args(argv)

    prs = Presentation(args.deck)
    slides = list(prs.slides)
    if not 1 <= args.slide <= len(slides):
        raise SystemExit(f"drop: в деке {len(slides)} слайдов, запрошен {args.slide}")
    slide = slides[args.slide - 1]
    shapes = list(slide.shapes)
    for i in args.shape:
        if not 0 <= i < len(shapes):
            raise SystemExit(f"drop: на слайде {args.slide} фигур {len(shapes)}, запрошена [{i}]")

    print(f"{args.deck}: слайд {args.slide}, фигур {len(shapes)}")
    for i in sorted(set(args.shape)):
        sh = shapes[i]
        bottom = ((sh.top or 0) + (sh.height or 0)) / _EMU_IN
        text = sh.text_frame.text.replace("\n", " | ")[:80] if sh.has_text_frame else ""
        print(f"  убрать [{i}] {sh.shape_type} name={sh.name!r} низ {bottom:.2f}in")
        if text:
            print(f"      {text!r}")
    if args.report:
        return 0

    tree = slide.shapes._spTree
    for i in sorted(set(args.shape), reverse=True):
        tree.remove(shapes[i]._element)

    out = args.out or args.deck
    prs.save(out)
    print(f"  осталось фигур {len(list(Presentation(out).slides)[args.slide - 1].shapes)}; "
          f"записано: {out}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
