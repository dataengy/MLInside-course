#!/usr/bin/env python3
"""Пересчитать часы в подстрочниках готовой .pptx — по фактическому порядку слайдов.

ЗАЧЕМ. Штатный путь — `just preza-notes apply content/…-content.yml`: часы считает
`preza_gen.notes` по контенту, из которого дека и собирается. Но у РУЧНОЙ линии деки
контент-YAML нет: файл правят в PowerPoint, слайды переставляют и вставляют, а шапка
`[MM:SS–MM:SS · ~N мин]` остаётся от прошлого порядка. Тогда часы врут — и врут молча:
метка `[~N мин]` в них верная, а абсолютное время уже нет.

Скрипт закрывает ровно этот разрыв: читает длительности прямо из заметок pptx, складывает
их по порядку слайдов и переставляет шапку. Формат шапки не изобретается — берётся
`preza_gen.notes.stamp`, тот же, что ставит генератор, поэтому обе линии остаются
сравнимыми, а прогон идемпотентен (`stamp` сам снимает предыдущую шапку).

ЧЕГО СКРИПТ НЕ ДЕЛАЕТ. Не придумывает длительность. Слайд без метки `[~N мин]` считается
нулевым и часов не получает — как `scope: stamped` в настройках сборки. Такие слайды
перечисляются в отчёте: если их много подряд, хронометраж лекции занижен, и это решается
проставлением меток, а не пересчётом.

ИСПОЛЬЗОВАНИЕ
    renumber_notes.py <deck.pptx> [--out OUT.pptx] [--set N=МИН]… [--report]

    --set N=МИН   задать длительность слайду N (1-based), если в заметке метки нет
    --out         куда писать; без него правится файл на месте
    --report      показать пересчёт и выйти, ничего не записывая
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "src"))

from pptx import Presentation  # noqa: E402

from preza_gen.notes import CLOCK_RE, _OLD_DURATION, clock, duration, stamp  # noqa: E402


def _head(start_min: float, end_min: float, mins: float) -> str:
    """Шапка заметки. Формат — тот же, что ставит генератор."""
    return f"[{clock(start_min)}–{clock(end_min)} · ~{mins:g} мин] "


# Формат шапки задан в preza_gen.notes и здесь только воспроизводится. Сверяемся с ним на
# импорте: если генератор сменит разделитель, скрипт обязан упасть сразу, а не тихо
# наставить деке шапок в устаревшем виде, которые потом никто не отличит от свежих.
_probe = stamp("[~1.5 мин] X", 10, 11.5)
assert _probe == _head(10, 11.5, 1.5) + "X", f"формат шапки разошёлся с preza_gen: {_probe!r}"


def _restamp(slide, head_new: str) -> bool:
    """Заменить ТОЛЬКО шапку в первом абзаце заметки, не трогая остальные раны.

    Почему не переписать заметку целиком: выделения в подстрочнике — не украшение, а разметка
    («!!самое важное!!», «**термин**», «~~второстепенное~~»), и живут они в отдельных ранах —
    на слайд их бывает под сорок. Перезапись абзацев одним плоским раном стёрла бы всю
    разметку разом, и заметить это можно было бы только глазами в PowerPoint.

    Шапка всегда стоит в начале первого абзаца, но может быть разрезана на раны посреди слова
    (``'[00:00–01:00 · ~1 '`` + ``'мин'`` + ``'] '``), поэтому съедаем раны по символам.
    """
    para = slide.notes_slide.notes_text_frame.paragraphs[0]
    runs = para.runs
    if not runs:
        return False
    joined = "".join(r.text for r in runs)
    # Длина старой шапки считается тем же порядком, каким её снимает preza_gen.notes.strip_head:
    # сперва часы (CLOCK_RE знает и нынешний формат, и прежний), затем отдельная метка
    # длительности старого вида. Иначе `[~2 мин]` осталась бы в тексте вторым экземпляром.
    lead = len(joined) - len(joined.lstrip())
    head_len = lead
    m = CLOCK_RE.match(joined[head_len:])
    head_len += m.end() if m else 0
    m = _OLD_DURATION.match(joined[head_len:])
    head_len += m.end() if m else 0

    eaten = 0
    for run in runs:
        if eaten >= head_len:
            break
        take = min(len(run.text), head_len - eaten)
        run.text = run.text[take:]
        eaten += take
    runs[0].text = head_new + runs[0].text
    return True


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("deck")
    ap.add_argument("--out")
    ap.add_argument("--set", action="append", default=[], metavar="N=МИН")
    ap.add_argument("--report", action="store_true")
    args = ap.parse_args(argv)

    overrides = {}
    for item in args.set:
        try:
            n, mins = item.split("=")
            overrides[int(n)] = float(mins.replace(",", "."))
        except ValueError:
            raise SystemExit(f"renumber: --set ждёт вид N=МИН, получено {item!r}")

    prs = Presentation(args.deck)
    slides = list(prs.slides)
    for n in overrides:
        if not 1 <= n <= len(slides):
            raise SystemExit(f"renumber: в деке {len(slides)} слайдов, задан {n}")

    cursor = 0.0
    changed = []
    untimed = []
    for i, slide in enumerate(slides, 1):
        text = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        if not text.strip():
            continue
        mins = overrides.get(i, duration(text))
        if mins is None:
            untimed.append(i)
            continue
        head_new = _head(cursor, cursor + mins, mins)
        if not text.lstrip().startswith(head_new):
            changed.append((i, cursor, cursor + mins, mins))
            if not args.report:
                _restamp(slide, head_new)
        cursor += mins

    print(f"{args.deck}: слайдов {len(slides)}, перештамповано {len(changed)}, "
          f"итог {clock(cursor)} ({cursor:g} мин)")
    if untimed:
        print(f"  без метки [~N мин] и потому без часов — {len(untimed)} шт.: "
              f"{untimed if len(untimed) <= 25 else str(untimed[:25]) + ' …'}")
    for i, a, b, m in changed[:8]:
        print(f"    #{i:>3}  [{clock(a)}–{clock(b)} · ~{m:g} мин]")
    if len(changed) > 8:
        print(f"    … ещё {len(changed) - 8}")
    if args.report:
        return 0

    out = args.out or args.deck
    prs.save(out)
    print(f"  записано: {out}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
