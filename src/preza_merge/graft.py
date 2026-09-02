"""preza_merge.graft — перенос слайдов между двумя ``.pptx`` с общим шаблоном.

ЗАЧЕМ. Лента `preza-merge` разбирает ФОРМАТИРОВАНИЕ чужого форка и превращает его в правила
генератора. Слайды она не переносит: правило «сделай кегль 21» вывести можно, а «возьми вот
этот слайд» — нельзя, его надо скопировать. Пока бэкенда не было, единственным ответом на
«человек дорисовал слайд в PowerPoint» оставалось «перенеси руками в content-YAML».

ГРАНИЦА, которую нельзя размывать. Результат графта — **рукописный артефакт**, а не
сгенерированная дека: из ``content/*-content.yml`` он не воспроизводится. Поэтому он и
кладётся в ``data/source/manual/``, и его содержимое всё равно придётся однажды перенести
в контент — иначе следующая сборка выдаст деку без этих слайдов. Графт покупает время,
а не отменяет перенос (см. docs/deck-manual-pass.md).

СОВМЕСТИМОСТЬ ФАЙЛОВ. Рассчитано на деки С ОБЩИМ ПРЕДКОМ: макет ищется по ИМЕНИ, и если
имени нет в целевой деке — падаем, а не подставляем молча другой макет (фигуры уехали бы
по месту). Тема, размер слайда и набор макетов не переносятся вообще.

НУМЕРАЦИЯ. Все номера 1-based и относятся к ИСХОДНОМУ состоянию файлов: порядок операций
внутри не влияет на то, как их считать. «С поправкой на уже вставленное» считать не нужно.
"""

from __future__ import annotations

import copy
import hashlib
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

#: Теги фигур, которые переносятся. Всё остальное (плейсхолдеры макета) не копируем.
_SHAPE_TAGS = (qn("p:sp"), qn("p:pic"), qn("p:graphicFrame"), qn("p:grpSp"), qn("p:cxnSp"))


class GraftError(RuntimeError):
    """Несовместимые файлы или невозможный план. Всегда громко, никогда молча."""


@dataclass(frozen=True)
class Op:
    """Одна операция плана.

    ``kind="insert"`` — вставить слайд ``src`` исходника ПОСЛЕ слайда ``at`` цели
    (``at=0`` — в начало). ``kind="replace"`` — заменить слайд ``at`` цели слайдом ``src``.
    """

    kind: str
    src: int
    at: int


def parse_ops(inserts: list[str], replaces: list[str]) -> list[Op]:
    """``["3:21"]`` → ``[Op("insert", 3, 21)]``. Разбор без обращения к файлам."""
    ops = []
    for raw, kind in [(x, "insert") for x in inserts] + [(x, "replace") for x in replaces]:
        try:
            left, right = raw.split(":")
            ops.append(Op(kind, int(left), int(right)))
        except ValueError:
            raise GraftError(f"--{kind} ждёт вид N:M, получено {raw!r}") from None
    return ops


def _validate(ops: list[Op], n_src: int, n_tgt: int) -> None:
    if not ops:
        raise GraftError("план пуст — нечего переносить; задайте --insert/--replace")
    for op in ops:
        if not 1 <= op.src <= n_src:
            raise GraftError(f"в исходнике {n_src} слайдов, запрошен {op.src}")
        if op.kind == "insert" and not 0 <= op.at <= n_tgt:
            raise GraftError(f"вставка после слайда {op.at}, а в цели их {n_tgt}")
        if op.kind == "replace" and not 1 <= op.at <= n_tgt:
            raise GraftError(f"замена слайда {op.at}, а в цели их {n_tgt}")


def _layout_by_name(prs, name: str):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    have = sorted({lay.name for m in prs.slide_masters for lay in m.slide_layouts})
    raise GraftError(f"макета {name!r} нет в целевой деке; есть: {have}")


def _media_index(part) -> dict[str, str]:
    """sha256 → rId уже связанных картинок: чтобы не плодить дубли в ``ppt/media/``."""
    idx = {}
    for rid, rel in part.rels.items():
        if rel.is_external:
            continue
        target = rel.target_part
        if str(target.partname).startswith("/ppt/media/"):
            idx[hashlib.sha256(target.blob).hexdigest()] = rid
    return idx


def _remap_media(src_slide, dst_slide, stats: dict) -> None:
    """Перевесить ``r:embed``/``r:link`` скопированных фигур на части целевой деки.

    ``added``/``reused`` считают СВЯЗИ слайда, а не файлы в пакете: связь с картинкой
    в OOXML своя у каждого слайда, и двум слайдам законно нужны две. Дедуп идёт по sha256
    содержимого и не даёт задвоиться самому файлу в ``ppt/media/`` — именно от этого
    зависит вес деки. Не читать ``added > 1`` как «картинка скопирована дважды».
    """
    dst_part = dst_slide.part
    known = _media_index(dst_part)
    for attr in (qn("r:embed"), qn("r:link")):
        for el in dst_slide.shapes._spTree.iter():
            rid = el.get(attr)
            if not rid:
                continue
            try:
                src_media = src_slide.part.related_part(rid)
            except KeyError:
                raise GraftError(f"в исходном слайде нет связи {rid} — файл повреждён") from None
            digest = hashlib.sha256(src_media.blob).hexdigest()
            if digest in known:
                stats["reused"] += 1
            else:
                known[digest] = dst_part.relate_to(src_media, src_slide.part.rels[rid].reltype)
                stats["added"] += 1
            el.set(attr, known[digest])


def _copy_slide(prs, src_slide):
    """Добавить в конец ``prs`` копию ``src_slide`` (фигуры + заметки) и вернуть её."""
    new = prs.slides.add_slide(_layout_by_name(prs, src_slide.slide_layout.name))

    # add_slide приносит плейсхолдеры макета — они не нужны, содержимое копируется целиком.
    spTree = new.shapes._spTree
    for shape in list(spTree):
        if shape.tag in _SHAPE_TAGS:
            spTree.remove(shape)
    for shape in src_slide.shapes._spTree:
        if shape.tag in _SHAPE_TAGS:
            spTree.append(copy.deepcopy(shape))

    if src_slide.has_notes_slide:
        dst_body = new.notes_slide.notes_text_frame._txBody
        for para in list(dst_body.findall(qn("a:p"))):
            dst_body.remove(para)
        for para in src_slide.notes_slide.notes_text_frame._txBody.findall(qn("a:p")):
            dst_body.append(copy.deepcopy(para))
    return new


def _move(prs, sld_id_el, index: int) -> None:
    lst = prs.slides._sldIdLst
    lst.remove(sld_id_el)
    lst.insert(index, sld_id_el)


def _drop(prs, index: int) -> None:
    lst = prs.slides._sldIdLst
    sld_id = lst[index]
    prs.part.drop_rel(sld_id.rId)
    lst.remove(sld_id)


def plan_report(target: Path, source: Path, ops: list[Op]) -> list[str]:
    """Человекочитаемый план, ничего не меняя. Для ``--report`` и для ревью до применения."""
    tgt, src = Presentation(str(target)), Presentation(str(source))
    src_slides = list(src.slides)
    _validate(ops, len(src_slides), len(tgt.slides._sldIdLst))

    def title(slide) -> str:
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                return sh.text_frame.text.strip().splitlines()[0][:64]
        return "(без текста)"

    lines = [
        f"target: {target}  ({len(tgt.slides._sldIdLst)} слайдов)",
        f"source: {source}  ({len(src_slides)} слайдов)",
    ]
    for op in ops:
        where = f"после #{op.at}" if op.kind == "insert" else f"вместо #{op.at}"
        lines.append(f"  {op.kind:<7} src#{op.src:>3} → {where:<12} {title(src_slides[op.src - 1])}")
    return lines


def graft(target: Path, source: Path, out: Path, ops: list[Op]) -> dict:
    """Применить план и записать результат в ``out``. Возвращает отчёт."""
    tgt, src = Presentation(str(target)), Presentation(str(source))
    src_slides = list(src.slides)
    _validate(ops, len(src_slides), len(tgt.slides._sldIdLst))

    stats = {"reused": 0, "added": 0}
    # Сначала копируем ВСЕ слайды в хвост, потом расставляем: так номера в плане не «плывут»
    # по ходу и остаются номерами исходного состояния.
    planned = [
        (op.at if op.kind == "insert" else op.at - 1, op.src,
         None if op.kind == "insert" else op.at)
        for op in ops
    ]
    made = []
    for at, s, drop_at in sorted(planned, key=lambda x: x[0]):
        new = _copy_slide(tgt, src_slides[s - 1])
        _remap_media(src_slides[s - 1], new, stats)
        made.append((at, tgt.slides._sldIdLst[-1], drop_at))

    # Расставляем от конца к началу: правки не сдвигают ещё не обработанные позиции.
    # Ключ включает порядковый номер, чтобы РАЗВЕРНУТЬ и группу с одинаковым `at`: каждая
    # следующая вставка в ту же позицию отодвигает предыдущую, и без разворота несколько
    # слайдов «после #21» легли бы задом наперёд.
    for _, (at, sld_id, _) in sorted(enumerate(made), key=lambda x: (x[1][0], x[0]), reverse=True):
        _move(tgt, sld_id, at)
    # Удаляем заменённые оригиналы: их индекс сдвинулся ровно на число вставок до него.
    for at, _, drop_at in sorted(made, key=lambda x: x[0], reverse=True):
        if drop_at is not None:
            shift = sum(1 for a2, _, _ in made if a2 <= at)
            _drop(tgt, drop_at - 1 + shift)

    out.parent.mkdir(parents=True, exist_ok=True)
    tgt.save(str(out))
    return {
        "out": out,
        "slides": len(Presentation(str(out)).slides._sldIdLst),
        "media_reused": stats["reused"],
        "media_added": stats["added"],
        "ops": len(ops),
    }
