"""preza_merge.model — a normalized, comparable model of a .pptx.

Only what the differ and the rule detectors are allowed to reason about: geometry in
inches, paragraph/run properties as raw OOXML attribute values, joined text, notes, theme
fonts and the master's body sizes. Picture bytes, animations and layout internals are out
of scope by design — the lane never rewrites a deck's content.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

_EMU_IN = 914400
_WS = re.compile(r"\s+")
# Paragraph/run/body attributes worth comparing — everything else is round-trip noise.
_BODY_KEYS = ("lIns", "tIns", "rIns", "bIns", "anchor", "wrap", "vert")
_PARA_KEYS = ("marL", "indent", "algn")
_RUN_KEYS = ("sz", "b", "i", "u")


@dataclass(frozen=True)
class Run:
    text: str
    size: int | None  # OOXML sz — hundredths of a point (2000 == 20 pt)
    bold: str | None
    italic: str | None
    underline: str | None
    font: str | None
    color: str | None


@dataclass(frozen=True)
class Para:
    level: int
    props: dict
    end_size: int | None  # endParaRPr@sz — survives when every run's size is cleared
    runs: list[Run]

    def text(self) -> str:
        """Runs joined, whitespace normalized — one paragraph, one string."""
        return _WS.sub(" ", "".join(r.text for r in self.runs)).strip()


@dataclass(frozen=True)
class Shape:
    name: str
    kind: str
    left: float
    top: float
    width: float
    height: float
    placeholder: str | None
    body_pr: dict
    paras: list[Para] = field(default_factory=list)
    table: list[list[str]] | None = None
    hyperlinks: list[str] = field(default_factory=list)
    line_color: str | None = None  # outline: hex, a theme role ("tx1"), or None when absent
    line_width: int | None = None  # EMU

    @property
    def bottom(self) -> float:
        return self.top + self.height

    def text(self) -> list[str]:
        """Non-empty paragraph texts, in order."""
        return [t for t in (p.text() for p in self.paras) if t]


@dataclass(frozen=True)
class Slide:
    n: int  # 1-based
    layout: str
    title: str
    notes: str
    shapes: list[Shape]
    shapes_title_name: str | None

    def by_name(self) -> dict[str, Shape]:
        return {s.name: s for s in self.shapes}


@dataclass(frozen=True)
class Deck:
    path: Path
    slides: list[Slide]
    theme_fonts: dict[str, str]
    master_body_sizes: dict[int, float]

    def titles(self) -> list[str]:
        return [s.title for s in self.slides]


def _in(value) -> float:
    return 0.0 if value is None else round(Emu(int(value)).inches, 3)


def _body_pr(text_frame) -> dict:
    node = text_frame._txBody.find(qn("a:bodyPr"))
    if node is None:
        return {}
    out = {k: node.get(k) for k in _BODY_KEYS if node.get(k) is not None}
    for child in node:
        tag = child.tag.split("}")[1]
        if tag in ("normAutofit", "spAutoFit", "noAutofit"):
            out["autofit"] = tag
    return out


def _para(paragraph) -> Para:
    props: dict = {}
    node = paragraph._p.find(qn("a:pPr"))
    end_size = None
    if node is not None:
        props = {k: node.get(k) for k in _PARA_KEYS if node.get(k) is not None}
        for tag in ("lnSpc", "spcBef", "spcAft"):
            el = node.find(qn(f"a:{tag}"))
            if el is not None and len(el):
                props[tag] = el[0].get("val")
        for tag in ("buNone", "buChar", "buAutoNum"):
            if node.find(qn(f"a:{tag}")) is not None:
                props[tag] = True
    end = paragraph._p.find(qn("a:endParaRPr"))
    if end is not None and end.get("sz"):
        end_size = int(end.get("sz"))
    return Para(
        level=paragraph.level,
        props=props,
        end_size=end_size,
        runs=[_run(r) for r in paragraph.runs],
    )


def _run(run) -> Run:
    node = run._r.find(qn("a:rPr"))
    size = font = color = bold = italic = underline = None
    if node is not None:
        size = int(node.get("sz")) if node.get("sz") else None
        bold, italic, underline = node.get("b"), node.get("i"), node.get("u")
        latin = node.find(qn("a:latin"))
        font = latin.get("typeface") if latin is not None else None
        fill = node.find(qn("a:solidFill"))
        if fill is not None and len(fill):
            color = fill[0].get("val") or fill[0].get("lastClr")
    return Run(run.text, size, bold, italic, underline, font, color)


def _hyperlinks(shape) -> list[str]:
    """External link targets inside a shape, in document order."""
    if not shape.has_text_frame:
        return []
    out = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            address = run.hyperlink.address
            if address:
                out.append(address)
    return out


def _line(shape) -> tuple[str | None, int | None]:
    """Outline colour + width of a shape, or (None, None) when it defines no line.

    The colour is returned RAW — a hex value or a theme role like "tx1" — because that is
    exactly what tells an accent border apart from the theme-text one the manager switched to.
    """
    props = shape._element.find(qn("p:spPr"))
    line = props.find(qn("a:ln")) if props is not None else None
    if line is None:
        return None, None
    fill = line.find(qn("a:solidFill"))
    color = None
    if fill is not None and len(fill):
        color = fill[0].get("val") or fill[0].get("lastClr")
    width = int(line.get("w")) if line.get("w") else None
    return color, width


def _shape(shape) -> Shape:
    paras: list[Para] = []
    body_pr: dict = {}
    if shape.has_text_frame:
        body_pr = _body_pr(shape.text_frame)
        paras = [_para(p) for p in shape.text_frame.paragraphs]
    table = None
    if getattr(shape, "has_table", False) and shape.has_table:
        table = [[c.text for c in row.cells] for row in shape.table.rows]
    placeholder = None
    if shape.is_placeholder:
        placeholder = f"{shape.placeholder_format.type}/{shape.placeholder_format.idx}"
    line_color, line_width = _line(shape)
    return Shape(
        name=shape.name,
        kind=str(shape.shape_type),
        left=_in(shape.left),
        top=_in(shape.top),
        width=_in(shape.width),
        height=_in(shape.height),
        placeholder=placeholder,
        body_pr=body_pr,
        paras=paras,
        table=table,
        hyperlinks=_hyperlinks(shape),
        line_color=line_color,
        line_width=line_width,
    )


def _theme_fonts(prs) -> dict[str, str]:
    for part in prs.part.package.iter_parts():
        if "theme" not in str(part.partname):
            continue
        xml = part.blob.decode("utf-8", "ignore")
        found = re.findall(r'<a:(major|minor)Font>\s*<a:latin typeface="([^"]*)"', xml)
        if found:
            return {k: v for k, v in found}
    return {}


def _master_body_sizes(prs) -> dict[int, float]:
    sizes: dict[int, float] = {}
    for master in prs.slide_masters:
        styles = master._element.find(qn("p:txStyles"))
        if styles is None:
            continue
        body = styles.find(qn("p:bodyStyle"))
        if body is None:
            continue
        for lvl, node in enumerate(body):
            props = node.find(qn("a:defRPr"))
            if props is not None and props.get("sz"):
                sizes[lvl] = int(props.get("sz")) / 100
        if sizes:
            break
    return sizes


def load(path: str | Path) -> Deck:
    """Normalize a .pptx into the comparable model."""
    p = Path(path).expanduser()
    prs = Presentation(str(p))
    slides = []
    for i, s in enumerate(prs.slides, 1):
        title_shape = s.shapes.title
        slides.append(
            Slide(
                n=i,
                layout=s.slide_layout.name,
                title=(title_shape.text if title_shape is not None else ""),
                notes=(s.notes_slide.notes_text_frame.text if s.has_notes_slide else ""),
                shapes=[_shape(sh) for sh in s.shapes],
                shapes_title_name=(title_shape.name if title_shape is not None else None),
            )
        )
    return Deck(
        path=p,
        slides=slides,
        theme_fonts=_theme_fonts(prs),
        master_body_sizes=_master_body_sizes(prs),
    )
