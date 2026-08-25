"""Synthetic decks for merge-lane tests — built on python-pptx's own template.

Deliberately independent of the course template (data/ is git-lfs): the merge lane must be
testable on any checkout.
"""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt


def _add(prs, title, bullets, *, sizes=None, body_width=None, body_top=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    body.name = "Text Placeholder 2"  # matches the real generator's shape-name prefix
    if body_width is not None:
        body.width = Inches(body_width)
    if body_top is not None:
        body.top = Inches(body_top)
    tf = body.text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        if sizes is not None:
            for r in p.runs:
                r.font.size = Pt(sizes)
    return slide


@pytest.fixture
def make_deck(tmp_path):
    """make_deck("name", [(title, [bullets]), ...], sizes=20) -> Path to a .pptx"""

    def _make(name: str, slides, *, sizes=None, body_width=None, body_top=None) -> Path:
        prs = Presentation()
        for title, bullets in slides:
            _add(prs, title, bullets, sizes=sizes, body_width=body_width, body_top=body_top)
        out = tmp_path / f"{name}.pptx"
        prs.save(str(out))
        return out

    return _make


@pytest.fixture
def add_shape():
    """add_shape(slide, name, top, height, left=1.0, width=2.0) -> Shape

    A plain autoshape renamed to match a detector's name prefix (e.g. "Rounded
    Rectangle 1", "Picture 1") — the geometry-based rules key off the name, not the
    PowerPoint shape type, so a rectangle stands in for a real code panel or picture.
    """
    from pptx.enum.shapes import MSO_SHAPE

    def _add_shape(slide, name, top, height, *, left=1.0, width=2.0):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.name = name
        return shape

    return _add_shape


@pytest.fixture
def add_table():
    """add_table(slide, name, top, left=1.0, width=2.0, height=1.0) -> GraphicFrame

    A 1x1 table renamed to match the "Table" prefix, at the given top — its cell
    content is irrelevant to R3, which only reads the shape's geometry.
    """

    def _add_table(slide, name, top, *, left=1.0, width=2.0, height=1.0):
        gframe = slide.shapes.add_table(
            1, 1, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        gframe.name = name
        return gframe

    return _add_table
