"""Applying a proposal: profile written, deck switched, patch built — and refusals."""

from pathlib import Path

import pytest
import yaml

from preza_merge import apply, graft, rules

_REPO = Path(__file__).resolve().parents[3]


def _formats(tmp_path) -> Path:
    src = yaml.safe_load((_REPO / "settings" / "formats.yml").read_text(encoding="utf-8"))
    p = tmp_path / "formats.yml"
    p.write_text(yaml.safe_dump(src, allow_unicode=True), encoding="utf-8")
    return p


def test_profile_inherits_the_base_and_overrides_accepted_keys(tmp_path):
    path = _formats(tmp_path)
    apply.write_profile(path, "merged", "classic", {"body_font": "inherit", "table_top": 2.45})

    data = yaml.safe_load(path.read_text(encoding="utf-8"))["formats"]["merged"]
    assert data["body_font"] == "inherit"
    assert data["table_top"] == 2.45
    assert data["bullets_width_narrow"] == 6.2  # inherited from classic
    assert set(data) == set(yaml.safe_load(path.read_text(encoding="utf-8"))["formats"]["classic"])


def test_existing_profiles_survive_a_rewrite(tmp_path):
    path = _formats(tmp_path)
    apply.write_profile(path, "merged", "classic", {"table_top": 2.45})
    apply.write_profile(path, "other", "classic", {"table_top": 3.0})
    formats = yaml.safe_load(path.read_text(encoding="utf-8"))["formats"]
    assert {"classic", "merged", "other"} <= set(formats)
    assert formats["merged"]["table_top"] == 2.45


def test_deck_format_is_set_surgically(tmp_path):
    content = tmp_path / "deck.yml"
    content.write_text(
        "deck:\n  out_name: X   # keep me\n  naming: increment\n  version_major: 3\n"
        "content:\n- kind: title\n  title: T\n",
        encoding="utf-8",
    )
    assert apply.set_deck_format(content, "merged") is True
    text = content.read_text(encoding="utf-8")
    assert "format: merged" in text
    assert "# keep me" in text  # the surgical edit must not reformat the file
    assert yaml.safe_load(text)["deck"]["format"] == "merged"


def test_deck_format_replaces_an_existing_value(tmp_path):
    content = tmp_path / "deck.yml"
    content.write_text(
        "deck:\n  out_name: X\n  format: classic\ncontent: []\n", encoding="utf-8"
    )
    apply.set_deck_format(content, "merged")
    assert yaml.safe_load(content.read_text(encoding="utf-8"))["deck"]["format"] == "merged"


def test_deck_format_ignores_a_lookalike_line_inside_a_code_panel(tmp_path):
    """A Data Engineering deck can plausibly hold `format: parquet` inside a slide's code
    block — the search must stay bounded to the `deck:` header, above `content:`, or an
    unscoped regex would silently rewrite the code sample instead."""
    content = tmp_path / "deck.yml"
    content.write_text(
        "deck:\n"
        "  out_name: X   # keep me\n"
        "  naming: increment\n"
        "  version_major: 3\n"
        "content:\n"
        "- kind: code\n"
        "  title: Хранение\n"
        "  code: |\n"
        "    df.write.format: parquet\n"
        "    df.save('/data/out')\n",
        encoding="utf-8",
    )
    before = content.read_text(encoding="utf-8")
    assert apply.set_deck_format(content, "merged") is True
    after = content.read_text(encoding="utf-8")

    doc = yaml.safe_load(after)
    assert doc["deck"]["format"] == "merged"
    # The code panel's own `format: parquet` line must survive byte-for-byte.
    assert "df.write.format: parquet" in after
    # Nothing changed except the inserted `format:` line under `deck:`.
    assert after == before.replace(
        "  out_name: X   # keep me\n", "  out_name: X   # keep me\n  format: merged\n"
    )


def test_apply_refuses_while_a_decision_is_missing(tmp_path):
    doc = {
        "proposal": {
            "deck": "content/x.yml",
            "base_pptx": "b.pptx",
            "ours_pptx": "o.pptx",
            "theirs_pptx": "t.pptx",
            "base_content_rev": "abc",
            "profile": "merged",
            "rules": [{"rule": "R1", "key": "body_font", "value": "inherit", "decision": None}],
            "regressions": [],
        }
    }
    cfg = rules.MergeConfig.load(_REPO / "settings" / "merge.yml")
    with pytest.raises(SystemExit, match="R1"):
        apply.run(
            doc,
            cfg,
            settings_yml=_REPO / "content" / "build_deck_v3-settings.yml",
            formats_path=_formats(tmp_path),
            patch_of="3.19",
            descr="x",
        )


def _graft_proposal(tmp_path):
    """Предложение, из которого бэкенд graft берёт только пару файлов."""
    from pptx import Presentation

    made = {}
    for role, titles in (("ours", ["A", "B"]), ("theirs", ["X"])):
        prs = Presentation()
        for title in titles:
            prs.slides.add_slide(prs.slide_layouts[1]).shapes.title.text = title
        made[role] = tmp_path / f"{role}.pptx"
        prs.save(str(made[role]))
    return {"proposal": {"rules": [], "regressions": [], "profile": "merged",
                         "deck": "content/x.yml",
                         "ours_pptx": str(made["ours"]),
                         "theirs_pptx": str(made["theirs"])}}


def test_graft_backend_needs_an_explicit_plan(tmp_path):
    """Предложение несёт правила ФОРМАТИРОВАНИЯ — плана слайдов из него не вывести.

    Поэтому пустой план обязан падать, а не тихо ничего не переносить: молчаливый no-op
    здесь выглядел бы как успешный перенос.
    """
    cfg = rules.MergeConfig.load(_REPO / "settings" / "merge.yml")
    with pytest.raises(graft.GraftError, match="план пуст"):
        apply.run(
            _graft_proposal(tmp_path), cfg,
            settings_yml=_REPO / "content" / "build_deck_v3-settings.yml",
            formats_path=_formats(tmp_path),
            patch_of="3.19", descr="x", backend="graft",
        )


def test_graft_backend_reports_a_missing_file_by_role(tmp_path):
    cfg = rules.MergeConfig.load(_REPO / "settings" / "merge.yml")
    prop = _graft_proposal(tmp_path)
    prop["proposal"]["theirs_pptx"] = str(tmp_path / "нет-такого.pptx")
    with pytest.raises(graft.GraftError, match="theirs_pptx"):
        apply.run(
            prop, cfg,
            settings_yml=_REPO / "content" / "build_deck_v3-settings.yml",
            formats_path=_formats(tmp_path),
            patch_of="3.19", descr="x", backend="graft", inserts=["1:1"],
        )


# ── комментарии ──────────────────────────────────────────────────────────────
#
# settings/formats.yml правится и машиной, и руками: числа профиля без объяснения,
# почему они именно такие, через месяц читаются как случайный набор. write_profile
# пишет round-trip и правит профиль ПО КЛЮЧАМ — оба свойства держатся тестами ниже,
# иначе первый же переход на safe_dump унесёт пояснения молча.
# Стратегии записи YAML и когда какая уместна — docs/glossary.md.

# Шапка файла каноническая: её дописывает код при каждой записи (как в publisher.plan_writer),
# поэтому здесь её нет — проверяем комментарии ВНУТРИ документа.
_COMMENTED_FORMATS = """\
formats:
  classic:
    # почему 6.51: подвал это логотип И строка «Материалы»
    visual_bottom: 6.51
    code_border: accent
"""


def _comment_lines(text: str) -> list[str]:
    return [ln.strip() for ln in text.splitlines() if ln.strip().startswith("#")]


def test_write_profile_keeps_comments(tmp_path):
    p = tmp_path / "formats.yml"
    p.write_text(_COMMENTED_FORMATS, encoding="utf-8")

    apply.write_profile(p, "alina", "classic", {"code_border": "none"})

    text = p.read_text(encoding="utf-8")
    assert "# почему 6.51: подвал это логотип И строка «Материалы»" in _comment_lines(text), (
        "комментарий внутри документа не пережил запись профиля"
    )
    assert text.count("named FORMATTING profiles") == 1, "каноническая шапка удвоилась"
    doc = yaml.safe_load(text)
    assert doc["formats"]["alina"]["code_border"] == "none"
    assert doc["formats"]["alina"]["visual_bottom"] == 6.51  # унаследовано от базы
    assert doc["formats"]["classic"]["code_border"] == "accent"  # база не тронута


def test_write_profile_is_byte_stable_on_repeat(tmp_path):
    """Второй прогон подряд обязан дать тот же файл — иначе каждая запись шумит диффом."""
    p = tmp_path / "formats.yml"
    p.write_text(_COMMENTED_FORMATS, encoding="utf-8")

    apply.write_profile(p, "alina", "classic", {"code_border": "none"})
    first = p.read_text(encoding="utf-8")
    apply.write_profile(p, "alina", "classic", {"code_border": "none"})
    assert p.read_text(encoding="utf-8") == first
