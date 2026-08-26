"""The normalized deck model — what the differ is allowed to see."""

from preza_merge import model


def test_model_captures_titles_bullets_and_geometry(make_deck):
    path = make_deck("a", [("Заголовок", ["раз", "два"])], sizes=20)
    deck = model.load(path)

    assert len(deck.slides) == 1
    slide = deck.slides[0]
    assert slide.title == "Заголовок"
    body = next(sh for sh in slide.shapes if sh.paras and sh.name != slide.shapes_title_name)
    assert body.text() == ["раз", "два"]
    assert body.width > 0
    assert {r.size for p in body.paras for r in p.runs} == {2000}


def test_text_joins_runs_and_normalizes_whitespace(make_deck):
    """A PowerPoint round-trip splits one bullet into many runs — the model must rejoin."""
    from pptx import Presentation

    path = make_deck("b", [("T", ["one"])])
    prs = Presentation(str(path))
    para = prs.slides[0].placeholders[1].text_frame.paragraphs[0]
    para.runs[0].text = "od"
    extra = para.add_run()
    extra.text = "in  \n one"
    prs.save(str(path))

    deck = model.load(path)
    slide = deck.slides[0]
    body = next(
        sh
        for sh in slide.shapes
        if sh.paras and sh.paras[0].runs and sh.name != slide.shapes_title_name
    )
    assert body.text() == ["odin one"]


def test_bottom_edge_is_derived(make_deck):
    path = make_deck("c", [("T", ["x"])])
    shape = next(sh for sh in model.load(path).slides[0].shapes if sh.height)
    assert abs(shape.bottom - (shape.top + shape.height)) < 1e-9


def test_shape_outline_is_captured(make_deck):
    """R11 lives here: a diff blind to outlines cannot notice a removed border."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    path = make_deck("e", [("T", ["x"])])
    prs = Presentation(str(path))
    shape = prs.slides[0].shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
    )
    shape.line.color.rgb = RGBColor.from_string("2419FF")
    shape.line.width = Pt(1)
    prs.save(str(path))

    found = next(sh for sh in model.load(path).slides[0].shapes if sh.name.startswith("Rounded"))
    assert found.line_color == "2419FF"
    assert found.line_width == 12700


def test_master_body_sizes_are_read(make_deck):
    path = make_deck("d", [("T", ["x"])])
    sizes = model.load(path).master_body_sizes
    assert sizes and all(v > 0 for v in sizes.values())
