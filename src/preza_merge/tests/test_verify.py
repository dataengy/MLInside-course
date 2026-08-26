"""Verification: residual geometry within tolerance, content invariants exact."""

import dataclasses
from pathlib import Path

import pytest

from preza_merge import model, rules, verify

_REPO = Path(__file__).resolve().parents[3]


@pytest.fixture
def cfg():
    return rules.MergeConfig.load(_REPO / "settings" / "merge.yml")


def test_small_residuals_pass(make_deck, cfg):
    # `top` is one of the TIGHT attributes (0.45″, see settings/merge.yml) — a real edge R2/R3
    # are meant to place exactly, unlike `width` which is deliberately loose for R4.
    a = model.load(make_deck("a", [("T", ["раз"])], body_top=1.0))
    b = model.load(make_deck("b", [("T", ["раз"])], body_top=1.2))  # 0.2" < 0.45" tolerance
    res = verify.structural(a, b, cfg)
    assert res.ok


def test_large_residuals_fail_with_a_slide_reference(make_deck, cfg):
    a = model.load(make_deck("a", [("T", ["раз"])], body_top=1.0))
    b = model.load(make_deck("b", [("T", ["раз"])], body_top=3.0))  # 2.0" > 0.45" tolerance
    res = verify.structural(a, b, cfg)
    assert not res.ok
    # Pin the actual slide reference — "1" alone would also match a digit inside a geometry
    # value and pass for the wrong reason.
    assert any("слайд 1" in m for m in res.mismatches)


def test_structural_requires_a_tolerance_for_every_geometry_attribute(make_deck, cfg):
    """A tolerances map missing one of diff._GEOM_ATTRS must raise, not silently skip it."""
    a = model.load(make_deck("a", [("T", ["раз"])], body_width=6.2))
    b = model.load(make_deck("b", [("T", ["раз"])], body_width=11.0))
    incomplete = dataclasses.replace(
        cfg, tolerances={k: v for k, v in cfg.tolerances.items() if k != "width"}
    )
    with pytest.raises(KeyError, match="width"):
        verify.structural(a, b, incomplete)


def test_invariants_catch_content_drift(make_deck):
    ours = model.load(make_deck("o", [("T", ["раз", "два"])]))
    merged_same = model.load(make_deck("m", [("T", ["раз", "два"])]))
    merged_drift = model.load(make_deck("d", [("T", ["раз", "ТРИ"])]))

    assert verify.invariants(ours, merged_same).ok
    bad = verify.invariants(ours, merged_drift)
    assert not bad.ok
    assert any("буллет" in m or "текст" in m for m in bad.mismatches)


def test_invariants_allow_uppercase_only_on_the_title_slide(make_deck):
    """R7 upcases ONLY slide 1 — the same case-only change on any other slide is real drift."""
    ours = model.load(make_deck("o", [("Алина", ["x"]), ("Бета", ["y"])]))
    merged_title_upcased = model.load(make_deck("m1", [("АЛИНА", ["x"]), ("Бета", ["y"])]))
    merged_other_upcased = model.load(make_deck("m2", [("Алина", ["x"]), ("БЕТА", ["y"])]))

    assert verify.invariants(ours, merged_title_upcased).ok

    bad = verify.invariants(ours, merged_other_upcased)
    assert not bad.ok
    assert any("слайд 2" in m for m in bad.mismatches)


def test_invariants_catch_a_lost_slide(make_deck):
    ours = model.load(make_deck("o", [("A", ["x"]), ("B", ["y"])]))
    merged = model.load(make_deck("m", [("A", ["x"])]))
    res = verify.invariants(ours, merged)
    assert not res.ok
    assert any("слайд" in m for m in res.mismatches)


def test_invariants_catch_lost_notes_and_links(make_deck):
    from pptx import Presentation

    ours_path = make_deck("o", [("A", ["x"])])
    prs = Presentation(str(ours_path))
    prs.slides[0].notes_slide.notes_text_frame.text = "заметка"
    run = prs.slides[0].placeholders[1].text_frame.paragraphs[0].runs[0]
    run.hyperlink.address = "https://example.com"
    prs.save(str(ours_path))

    ours = model.load(ours_path)
    merged = model.load(make_deck("m", [("A", ["x"])]))
    res = verify.invariants(ours, merged)
    assert not res.ok
    assert any("ссыл" in m for m in res.mismatches)
    assert any("заметк" in m for m in res.mismatches)
