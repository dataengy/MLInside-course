"""Pairwise diff — what counts as a real change and what is round-trip noise."""

from preza_merge import diff, model


def test_identical_decks_report_nothing(make_deck):
    a = make_deck("a", [("T", ["раз", "два"])], sizes=20)
    b = make_deck("b", [("T", ["раз", "два"])], sizes=20)
    rep = diff.compare(model.load(a), model.load(b))
    assert rep.counts["text"] == 0
    assert rep.counts["geometry"] == 0
    assert rep.counts["font"] == 0


def test_run_splitting_is_not_a_text_change(make_deck):
    """Regression: a Google Slides round-trip split runs and looked like 34 text edits."""
    from pptx import Presentation

    a = make_deck("a", [("T", ["раз два"])])
    b = make_deck("b", [("T", ["раз два"])])
    prs = Presentation(str(b))
    para = prs.slides[0].placeholders[1].text_frame.paragraphs[0]
    para.runs[0].text = "раз"
    para.add_run().text = " два"
    prs.save(str(b))

    rep = diff.compare(model.load(a), model.load(b))
    assert rep.counts["text"] == 0


def test_cleared_run_sizes_are_counted(make_deck):
    a = make_deck("a", [("T", ["раз", "два"])], sizes=20)
    b = make_deck("b", [("T", ["раз", "два"])])
    rep = diff.compare(model.load(a), model.load(b))
    assert rep.slides[0].runs_size_cleared == 2
    assert rep.counts["font"] >= 1


def test_geometry_changes_are_listed_per_shape(make_deck):
    a = make_deck("a", [("T", ["раз"])], body_width=6.2)
    b = make_deck("b", [("T", ["раз"])], body_width=10.0)
    rep = diff.compare(model.load(a), model.load(b))
    changes = [c for c in rep.slides[0].geometry if c.attr == "width"]
    assert changes and abs(changes[0].after - 10.0) < 0.02


def test_lost_notes_and_merged_paragraphs_are_flagged(make_deck):
    from pptx import Presentation

    a = make_deck("a", [("T", ["раз", "два"])])
    prs = Presentation(str(a))
    prs.slides[0].notes_slide.notes_text_frame.text = "заметка"
    prs.save(str(a))

    b = make_deck("b", [("T", ["раз два"])])
    rep = diff.compare(model.load(a), model.load(b))
    assert rep.slides[0].notes_lost is True
    assert rep.slides[0].paras_before > rep.slides[0].paras_after


def test_slide_count_mismatch_compares_the_common_prefix(make_deck):
    a = make_deck("a", [("T1", ["x"]), ("T2", ["y"])])
    b = make_deck("b", [("T1", ["x"])])
    rep = diff.compare(model.load(a), model.load(b))
    assert rep.slide_count == (2, 1)
    assert len(rep.slides) == 1
