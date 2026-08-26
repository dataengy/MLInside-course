"""Rule detectors: a systematic change becomes a profile key, a one-off does not."""

from pathlib import Path

import pytest
from pptx import Presentation

from preza_merge import diff, model, rules

_REPO = Path(__file__).resolve().parents[3]


@pytest.fixture
def cfg():
    return rules.MergeConfig.load(_REPO / "settings" / "merge.yml")


def _found(findings, rule):
    return next((f for f in findings if f.rule == rule), None)


def test_r1_fires_when_explicit_sizes_are_cleared(make_deck, cfg):
    base = model.load(make_deck("b", [("A", ["раз", "два"]), ("B", ["три"])], sizes=20))
    theirs = model.load(make_deck("t", [("A", ["раз", "два"]), ("B", ["три"])]))
    found = _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R1")
    assert found is not None
    assert found.key == "body_font" and found.value == "inherit"
    assert "3" in found.evidence  # three runs lost their explicit size


def test_r1_stays_silent_below_the_share_threshold(make_deck, cfg):
    """One slide out of five is a one-off, not a rule."""
    slides = [(f"S{i}", ["раз"]) for i in range(5)]
    base = model.load(make_deck("b", slides, sizes=20))
    path = make_deck("t", slides, sizes=20)
    prs = Presentation(str(path))
    for run in prs.slides[0].placeholders[1].text_frame.paragraphs[0].runs:
        run.font.size = None
    prs.save(str(path))
    theirs = model.load(path)
    assert _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R1") is None


def test_r2_fires_on_a_genuine_bottom_anchored_fork(make_deck, add_shape, cfg):
    """Base panels sit at varied bottoms; theirs anchors all of them to the same edge."""
    slides = [(f"S{i}", ["x"]) for i in range(4)]
    base_path = make_deck("b", slides)
    prs = Presentation(str(base_path))
    for i, slide in enumerate(prs.slides):
        add_shape(slide, f"Rounded Rectangle {i}", top=1.0 + i * 0.5, height=2.0 + i * 0.5)
    prs.save(str(base_path))
    base = model.load(base_path)

    theirs_path = make_deck("t", slides)
    prs = Presentation(str(theirs_path))
    for i, slide in enumerate(prs.slides):
        add_shape(slide, f"Rounded Rectangle {i}", top=5.0, height=2.0)  # bottom = 7.0 everywhere
    prs.save(str(theirs_path))
    theirs = model.load(theirs_path)

    found = _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R2")
    assert found is not None
    assert found.value == {"visual_anchor": "bottom", "visual_bottom": 7.0}
    assert found.share == 1.0


def test_r2_is_not_fooled_by_a_single_outlier(make_deck, add_shape, cfg):
    """3 of 4 panels sit at the same bottom edge in theirs; 1 moved wildly — the share
    gate (4/4 required, one is an outlier) must reject it."""
    slides = [(f"S{i}", ["x"]) for i in range(4)]
    base_path = make_deck("b", slides)
    prs = Presentation(str(base_path))
    for i, slide in enumerate(prs.slides):
        add_shape(slide, f"Rounded Rectangle {i}", top=3.0, height=2.0)  # bottom = 5.0
    prs.save(str(base_path))
    base = model.load(base_path)

    theirs_path = make_deck("t", slides)
    prs = Presentation(str(theirs_path))
    for i, slide in enumerate(prs.slides):
        top = 7.0 if i == 3 else 3.0  # slide 3's panel jumps to bottom = 9.0
        add_shape(slide, f"Rounded Rectangle {i}", top=top, height=2.0)
    prs.save(str(theirs_path))
    theirs = model.load(theirs_path)

    assert _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R2") is None


def test_r2_declines_when_nothing_moved(make_deck, add_shape, cfg):
    """Base and theirs place every panel at the same bottom edge — R2 must decline
    even though the (unconditional) share of aligned panels is 1.0."""
    slides = [(f"S{i}", ["x"]) for i in range(4)]
    base_path = make_deck("b", slides)
    prs = Presentation(str(base_path))
    for i, slide in enumerate(prs.slides):
        add_shape(slide, f"Rounded Rectangle {i}", top=3.0, height=2.0)
    prs.save(str(base_path))
    base = model.load(base_path)

    theirs_path = make_deck("t", slides)
    prs = Presentation(str(theirs_path))
    for i, slide in enumerate(prs.slides):
        add_shape(slide, f"Rounded Rectangle {i}", top=3.0, height=2.0)
    prs.save(str(theirs_path))
    theirs = model.load(theirs_path)

    assert _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R2") is None


def test_r3_fires_when_tables_move_down_by_a_consistent_amount(make_deck, add_table, cfg):
    base_path = make_deck("b", [(f"S{i}", ["x"]) for i in range(3)])
    prs = Presentation(str(base_path))
    for i, slide in enumerate(prs.slides):
        add_table(slide, f"Table {i}", top=1.5)
    prs.save(str(base_path))
    base = model.load(base_path)

    theirs_path = make_deck("t", [(f"S{i}", ["x"]) for i in range(3)])
    prs = Presentation(str(theirs_path))
    for i, slide in enumerate(prs.slides):
        add_table(slide, f"Table {i}", top=2.47)
    prs.save(str(theirs_path))
    theirs = model.load(theirs_path)

    found = _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R3")
    assert found is not None
    assert found.key == "table_top" and found.value == 2.47
    assert found.share == 1.0


def test_r3_stays_silent_when_tables_did_not_move(make_deck, add_table, cfg):
    slides = [(f"S{i}", ["x"]) for i in range(3)]
    base_path = make_deck("b", slides)
    prs = Presentation(str(base_path))
    for i, slide in enumerate(prs.slides):
        add_table(slide, f"Table {i}", top=1.5)
    prs.save(str(base_path))
    base = model.load(base_path)

    theirs_path = make_deck("t", slides)
    prs = Presentation(str(theirs_path))
    for i, slide in enumerate(prs.slides):
        add_table(slide, f"Table {i}", top=1.5)
    prs.save(str(theirs_path))
    theirs = model.load(theirs_path)

    assert _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R3") is None


def test_r6_fires_with_share_reflecting_the_eligible_cohort(make_deck, cfg):
    """5 slides start with an empty body placeholder (eligible); the reviewer's tool
    physically removed it on 4 of the 5 — share must read 4/5 against that cohort."""
    slides = [(f"S{i}", []) for i in range(5)]  # no bullets => the body stays empty
    base = model.load(make_deck("b", slides))

    theirs_path = make_deck("t", slides)
    prs = Presentation(str(theirs_path))
    for slide in list(prs.slides)[:4]:
        body = next(sh for sh in slide.shapes if sh.name == "Text Placeholder 2")
        body._element.getparent().remove(body._element)
    prs.save(str(theirs_path))
    theirs = model.load(theirs_path)

    found = _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R6")
    assert found is not None
    assert found.key == "drop_empty_placeholders" and found.value is True
    assert found.share == 0.8


def test_r4_fires_when_the_bullet_column_widens(make_deck, cfg):
    """A code panel on the slide is what makes it eligible for R4 at all."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    def _with_panel(name, width):
        path = make_deck(name, [("A", ["x"]), ("B", ["y"])], body_width=width)
        prs = Presentation(str(path))
        for slide in prs.slides:
            slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
            )
        prs.save(str(path))
        return model.load(path)

    base = _with_panel("b", 6.2)
    theirs = _with_panel("t", 10.0)
    found = _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R4")
    assert found is not None and found.key == "bullets_width" and found.value == "adaptive"


def test_r4_stays_silent_when_no_slide_carries_a_visual(make_deck, cfg):
    """Bullets can widen all they like; without a code panel or picture on the slide R4's
    cohort is empty — the rule must not fire on a bullets-only slide."""
    base = model.load(make_deck("b", [("A", ["x"]), ("B", ["y"])], body_width=6.2))
    theirs = model.load(make_deck("t", [("A", ["x"]), ("B", ["y"])], body_width=10.0))
    assert _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R4") is None


def test_r1_denominator_is_slides_with_explicit_sizes_not_all_paragraphs(make_deck, cfg):
    """4 of 8 slides start with explicit run sizes; the other 4 never had one. Clearing all
    4 sized slides must read as share=1.0 (4/4), not 0.5 (4/8) — the unsized slides, which
    would dilute an "any paragraph" cohort, must not enter R1's denominator at all."""
    sized = [(f"S{i}", ["раз"]) for i in range(4)]
    unsized = [(f"U{i}", ["два"]) for i in range(4)]

    base_path = make_deck("b", sized, sizes=20)
    prs = Presentation(str(base_path))
    for title, bullets in unsized:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1]
        body.name = "Text Placeholder 2"
        body.text_frame.text = bullets[0]  # no explicit run size set
    prs.save(str(base_path))
    base = model.load(base_path)

    # theirs never had explicit sizes anywhere — every sized run in base gets cleared.
    theirs = model.load(make_deck("t", sized + unsized))

    found = _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R1")
    assert found is not None
    assert found.share == 1.0


def test_merge_config_threshold_uses_override_or_falls_back_to_min_share(cfg):
    assert cfg.threshold("R4") == 0.75
    assert cfg.threshold("R6") == 0.70
    assert cfg.threshold("R1") == cfg.min_share


def test_merge_config_threshold_raises_on_unknown_rule_id(cfg):
    """A typo'd rule id must not silently fall back to min_share."""
    with pytest.raises(ValueError):
        cfg.threshold("R99")


def test_merge_config_load_raises_on_a_bogus_override_rule_id(tmp_path):
    """A typo in settings/merge.yml's min_share_overrides must fail loud at load time,
    not go unnoticed because it never happens to match a real rule id."""
    bad = tmp_path / "bad_merge.yml"
    bad.write_text(
        """
merge:
  min_share: 0.8
  min_share_overrides:
    R99: 0.5
  tolerances: {left: 0.4, top: 0.4, width: 0.4, height: 0.4}
  report_dir: docs/reviews/merge
  fork_markers: [" (1)"]
  fork_search_dir: ~/Downloads
  default_profile: merged
  base_profile: classic
""",
        encoding="utf-8",
    )
    with pytest.raises(ValueError):
        rules.MergeConfig.load(bad)


def test_r11_fires_when_the_panel_border_goes_dark(make_deck, cfg):
    """The manager removed the blue outline on every code panel."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    def _with_panels(name, color):
        path = make_deck(name, [("A", ["x"]), ("B", ["y"])])
        prs = Presentation(str(path))
        for slide in prs.slides:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
            )
            shape.line.color.rgb = RGBColor.from_string(color)
        prs.save(str(path))
        return model.load(path)

    base = _with_panels("b", "2419FF")
    theirs = _with_panels("t", "1A1A1A")
    found = _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R11")
    assert found is not None
    assert found.key == "code_border" and found.value == "dark"


def test_r8_regression_on_theme_font_swap(make_deck, cfg):
    """A theme-font swap is an export artefact — reported, never merged."""
    base = model.load(make_deck("b", [("A", ["x"])]))
    loaded = model.load(make_deck("t", [("A", ["x"])]))
    theirs = model.Deck(
        path=loaded.path,
        slides=loaded.slides,
        theme_fonts={"major": "Calibri Light", "minor": "Calibri"},
        master_body_sizes=loaded.master_body_sizes,
    )
    found = _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R8")
    assert found is not None and found.kind == "regression" and found.key is None


def test_r10_regression_when_notes_are_lost(make_deck, cfg):
    base_path = make_deck("b", [("A", ["x"])])
    prs = Presentation(str(base_path))
    prs.slides[0].notes_slide.notes_text_frame.text = "заметка"
    prs.save(str(base_path))
    base = model.load(base_path)
    theirs = model.load(make_deck("t", [("A", ["x"])]))
    found = _found(rules.detect(base, theirs, diff.compare(base, theirs), cfg), "R10")
    assert found is not None and found.kind == "regression"


def test_config_loads_thresholds_from_yaml(cfg):
    assert cfg.min_share == 0.8
    # left/top/height are the TIGHT edges (R2/R3 place them exactly); width is deliberately
    # the loosest — R4 picks BINARY between a narrow and a full column (see settings/merge.yml).
    assert cfg.tolerances["left"] == cfg.tolerances["top"] == cfg.tolerances["height"] == 0.45
    assert cfg.tolerances["width"] == 5.5
    assert cfg.min_share_overrides["R4"] == 0.75
    assert cfg.min_share_overrides["R6"] == 0.70
