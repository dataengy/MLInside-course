"""graft: слайды переносятся с заметками и картинками, план проверяется, номера не «плывут».

Деки синтетические (conftest строит их на шаблоне python-pptx) — лента обязана быть
тестируемой на любом чекауте, без git-lfs.
"""

import pytest
from pptx import Presentation
from pptx.util import Inches

from preza_merge import graft as g


def _titles(path):
    return [
        next((sh.text_frame.text for sh in s.shapes if sh.has_text_frame and sh.text_frame.text), "")
        for s in Presentation(str(path)).slides
    ]


def _deck(path, titles, *, notes=None, picture=None):
    prs = Presentation()
    for i, t in enumerate(titles):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = t
        if notes and i in notes:
            slide.notes_slide.notes_text_frame.text = notes[i]
        if picture and i in picture:
            slide.shapes.add_picture(str(picture[i]), Inches(1), Inches(1), Inches(1))
    prs.save(str(path))
    return path


@pytest.fixture()
def png(tmp_path):
    """Однопиксельный PNG — минимальная валидная картинка, без внешних файлов."""
    import base64

    p = tmp_path / "dot.png"
    p.write_bytes(base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
    ))
    return p


def test_insert_puts_slide_after_the_named_one(tmp_path):
    tgt = _deck(tmp_path / "t.pptx", ["A", "B", "C"])
    src = _deck(tmp_path / "s.pptx", ["X", "Y"])
    out = tmp_path / "o.pptx"

    g.graft(tgt, src, out, [g.Op("insert", 2, 1)])  # Y после A

    assert _titles(out) == ["A", "Y", "B", "C"]


def test_replace_swaps_the_target_slide(tmp_path):
    tgt = _deck(tmp_path / "t.pptx", ["A", "B", "C"])
    src = _deck(tmp_path / "s.pptx", ["X"])
    out = tmp_path / "o.pptx"

    g.graft(tgt, src, out, [g.Op("replace", 1, 2)])  # X вместо B

    assert _titles(out) == ["A", "X", "C"]


def test_numbers_refer_to_the_original_state(tmp_path):
    """Несколько операций разом: номера считаются от исходных файлов, а не «с поправкой»."""
    tgt = _deck(tmp_path / "t.pptx", ["A", "B", "C"])
    src = _deck(tmp_path / "s.pptx", ["X", "Y", "Z"])
    out = tmp_path / "o.pptx"

    g.graft(tgt, src, out, [
        g.Op("insert", 1, 0),      # X в начало
        g.Op("insert", 2, 3),      # Y после C
        g.Op("replace", 3, 2),     # Z вместо B
    ])

    assert _titles(out) == ["X", "A", "Z", "C", "Y"]


def test_two_inserts_after_the_same_slide_keep_their_order(tmp_path):
    tgt = _deck(tmp_path / "t.pptx", ["A", "B"])
    src = _deck(tmp_path / "s.pptx", ["X", "Y"])
    out = tmp_path / "o.pptx"

    g.graft(tgt, src, out, [g.Op("insert", 1, 1), g.Op("insert", 2, 1)])

    assert _titles(out) == ["A", "X", "Y", "B"]


def test_notes_travel_with_the_slide(tmp_path):
    tgt = _deck(tmp_path / "t.pptx", ["A"])
    src = _deck(tmp_path / "s.pptx", ["X"], notes={0: "подстрочник докладчика"})
    out = tmp_path / "o.pptx"

    g.graft(tgt, src, out, [g.Op("insert", 1, 1)])

    notes = Presentation(str(out)).slides[1].notes_slide.notes_text_frame.text
    assert "подстрочник докладчика" in notes


def test_identical_pictures_are_not_duplicated_in_the_package(tmp_path, png):
    """Одна картинка на двух перенесённых слайдах даёт ОДНУ часть ppt/media в пакете.

    Считать по `media_added` нельзя: связь с картинкой в OOXML своя у каждого слайда, и
    двум слайдам законно нужны две связи. Дублируется или нет САМ ФАЙЛ — видно только
    в списке частей пакета, и проверять надо именно его: от этого зависит вес деки.
    """
    import zipfile

    tgt = _deck(tmp_path / "t.pptx", ["A"])
    src = _deck(tmp_path / "s.pptx", ["X", "Y"], picture={0: png, 1: png})
    out = tmp_path / "o.pptx"

    g.graft(tgt, src, out, [g.Op("insert", 1, 1), g.Op("insert", 2, 1)])

    parts = [n for n in zipfile.ZipFile(out).namelist() if n.startswith("ppt/media/")]
    assert len(parts) == 1, f"файл картинки задвоился в пакете: {parts}"


def test_empty_plan_is_an_error_not_a_no_op(tmp_path):
    tgt = _deck(tmp_path / "t.pptx", ["A"])
    src = _deck(tmp_path / "s.pptx", ["X"])
    with pytest.raises(g.GraftError, match="план пуст"):
        g.graft(tgt, src, tmp_path / "o.pptx", [])


@pytest.mark.parametrize(
    ("op", "match"),
    [
        (g.Op("insert", 9, 1), "в исходнике"),
        (g.Op("insert", 1, 9), "в цели"),
        (g.Op("replace", 1, 9), "в цели"),
    ],
)
def test_out_of_range_fails_loudly(tmp_path, op, match):
    tgt = _deck(tmp_path / "t.pptx", ["A"])
    src = _deck(tmp_path / "s.pptx", ["X"])
    with pytest.raises(g.GraftError, match=match):
        g.graft(tgt, src, tmp_path / "o.pptx", [op])


def test_missing_layout_name_fails_instead_of_guessing(tmp_path):
    """Макета нет в цели — падаем: подставить чужой значит молча сдвинуть все фигуры."""
    tgt = _deck(tmp_path / "t.pptx", ["A"])
    src = _deck(tmp_path / "s.pptx", ["X"])
    prs = Presentation(str(tgt))
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == "Title and Content":
            layout.name = "Совсем другое имя"
    prs.save(str(tgt))

    with pytest.raises(g.GraftError, match="макета"):
        g.graft(tgt, src, tmp_path / "o.pptx", [g.Op("insert", 1, 1)])


def test_parse_ops_rejects_malformed_pairs():
    with pytest.raises(g.GraftError, match="ждёт вид N:M"):
        g.parse_ops(["3"], [])


def test_plan_report_changes_nothing(tmp_path):
    tgt = _deck(tmp_path / "t.pptx", ["A", "B"])
    src = _deck(tmp_path / "s.pptx", ["X"])
    before = tgt.read_bytes()

    lines = g.plan_report(tgt, src, [g.Op("insert", 1, 1)])

    assert tgt.read_bytes() == before
    assert any("insert" in ln and "X" in ln for ln in lines)
