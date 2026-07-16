"""preza_gen.renderers.pptx — render a Content model onto the MLInside template → .pptx.

Refactor of the original single-file generator: layout map + slide builders driven by SlideSpec.
Reuses the source deck's media (ppt/media/*). Output is hardlinked to Config.downloads_link so the
publish Stop-hook still fires; provenance is stamped into the file's core-properties.
"""

from __future__ import annotations

import contextlib
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from ..settings import Config, Content, ImageBox
from ..utils import fit_image_box, hardlink_or_copy, log, parse_runs


def ensure_media(cfg: Config) -> Path:
    if cfg.media_dir and Path(cfg.media_dir).is_dir() and any(Path(cfg.media_dir).iterdir()):
        return Path(cfg.media_dir)
    media = cfg.out_dir.parent / "source" / "media"
    media.mkdir(parents=True, exist_ok=True)
    if not any(media.iterdir()):
        with zipfile.ZipFile(str(cfg.source_deck)) as z:
            for n in z.namelist():
                if n.startswith("ppt/media/"):
                    (media / Path(n).name).write_bytes(z.read(n))
        log.info(f"extracted media → {media}")
    return media


def _clear(prs) -> None:
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        prs.part.drop_rel(sid.get(qn("r:id")))
        lst.remove(sid)


def _theme(cfg: Config) -> dict:
    t = cfg.theme
    return {
        "accent": RGBColor.from_string(t.accent),
        "white": RGBColor.from_string(t.white),
        "dark": RGBColor.from_string(t.dark),
        "alt": RGBColor.from_string(t.alt),
        "font": t.font,
    }


def _set_body(slide, bullets, sizes: dict | None = None) -> None:
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(bullets):
        text, lvl = (item[0], item[1]) if isinstance(item, (list, tuple)) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = lvl
        if sizes:
            pt = sizes.get(lvl, sizes.get(max(sizes)))
            for r in p.runs:
                r.font.size = Pt(pt)


def _add_pic(slide, path: str, box: ImageBox) -> None:
    bl, bt = Inches(box.left), Inches(box.top)
    bw, bh = int(Inches(box.width)), int(Inches(box.height))
    iw, ih = Image.open(path).size
    w, h = fit_image_box(iw, ih, bw, bh)
    slide.shapes.add_picture(path, bl + (bw - w) // 2, bt + (bh - h) // 2, width=w, height=h)


def _cell(cell, text, *, bold, color, size, fill, font) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.text = text
    for r in p.runs:
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def _add_table(slide, table: dict, theme: dict) -> None:
    headers, rows = table["headers"], table["rows"]
    ratios = table.get("col_ratios")
    left, top, width = Inches(0.45), Inches(1.5), Inches(12.43)
    tbl = slide.shapes.add_table(
        len(rows) + 1, len(headers), left, top, width, Inches(0.45 + 0.4 * len(rows))
    ).table
    tbl.first_row = True
    tbl.horz_banding = False
    if ratios:
        tot = sum(ratios)
        for i, rr in enumerate(ratios):
            tbl.columns[i].width = Emu(int(int(width) * rr / tot))
    for c, h in enumerate(headers):
        _cell(
            tbl.cell(0, c),
            h,
            bold=True,
            color=theme["white"],
            size=13,
            fill=theme["accent"],
            font=theme["font"],
        )
    for ri, row in enumerate(rows, 1):
        for c, v in enumerate(row):
            _cell(
                tbl.cell(ri, c),
                v,
                bold=False,
                color=theme["dark"],
                size=12,
                fill=(theme["white"] if ri % 2 else theme["alt"]),
                font=theme["font"],
            )


def _add_code(slide, code: str, caption: str, theme: dict, box: ImageBox, style: dict) -> None:
    """Render a code snippet in a dark rounded monospace panel (optional caption above it)."""
    if caption:
        cap = slide.shapes.add_textbox(
            Inches(box.left), Inches(box.top - 0.34), Inches(box.width), Inches(0.3)
        )
        cr = cap.text_frame.paragraphs[0].add_run()
        cr.text = caption
        cr.font.size = Pt(12)
        cr.font.bold = True
        cr.font.name = theme["font"]
        cr.font.color.rgb = theme["dark"]
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(box.left),
        Inches(box.top),
        Inches(box.width),
        Inches(box.height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(style["bg"])
    shp.line.color.rgb = theme["accent"]
    shp.line.width = Pt(1)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.14)
    tf.margin_top = tf.margin_bottom = Inches(0.1)
    fg = RGBColor.from_string(style["fg"])
    for i, line in enumerate(code.rstrip("\n").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT   # auto-shapes default to centered — force left for code
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = style["font"]
        r.font.size = Pt(style["size"])
        r.font.color.rgb = fg


def _set_notes(slide, text: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    for i, para in enumerate(text.split("\n\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        for seg, b, u, it in parse_runs(para):
            r = p.add_run()
            r.text = seg
            if b:
                r.font.bold = True
            if u:
                r.font.underline = True
            if it:
                r.font.italic = True


def _add_materials(slide, materials: list, theme: dict) -> None:
    if not materials:
        return
    # bottom-right corner, right-aligned (clear of the MLINSIDE logo bottom-left)
    box = slide.shapes.add_textbox(Inches(4.75), Inches(6.98), Inches(8.3), Inches(0.42))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    lead = p.add_run()
    lead.text = "📚 Материалы: "
    lead.font.size = Pt(9)
    lead.font.bold = True
    lead.font.name = theme["font"]
    lead.font.color.rgb = theme["dark"]
    for i, m in enumerate(materials):
        if i:
            sep = p.add_run()
            sep.text = "  ·  "
            sep.font.size = Pt(9)
            sep.font.color.rgb = theme["dark"]
        r = p.add_run()
        r.text = m.get("label") or m["url"]
        r.font.size = Pt(9)
        r.font.name = theme["font"]
        r.hyperlink.address = m["url"]


def _provenance(prs, cfg: Config) -> None:
    try:
        cp = prs.core_properties
        cp.comments = (
            f"preza_gen v3 · template={cfg.template.name} · source={cfg.source_deck.name} · "
            f"media reused from source · hardlinked to {cfg.downloads_link}"
        )
        cp.keywords = "preza_gen;v3;MLInside;hardlink"
    except Exception:  # pragma: no cover
        pass


def render(cfg: Config, content: Content) -> tuple[Path, str | None]:
    """Build the .pptx from Content; return (output_path, hardlink_method)."""
    media = ensure_media(cfg)
    prs = Presentation(str(cfg.template))
    _clear(prs)
    layouts = {lay.name: lay for m in prs.slide_masters for lay in m.slide_layouts}
    theme = _theme(cfg)

    for spec in content.slides:
        s = prs.slides.add_slide(layouts[cfg.layouts[spec.kind]])
        s.shapes.title.text = spec.title
        if spec.kind == "title" and spec.subtitle:
            with contextlib.suppress(Exception):
                s.placeholders[1].text = spec.subtitle
        elif spec.kind == "table" and spec.table:
            _add_table(s, spec.table, theme)
        elif spec.kind in ("agenda", "content"):
            side = bool(spec.image or (spec.code and spec.bullets))
            if side:
                ph = s.placeholders[1]
                l0, t0, h0 = ph.left, ph.top, ph.height
                ph.left, ph.top, ph.height, ph.width = l0, t0, h0, Inches(6.2)
            if spec.bullets:
                _set_body(s, spec.bullets, cfg.body_font["with_image" if side else "bullets_only"])
            else:
                s.placeholders[1].text_frame.clear()
            if spec.image:
                _add_pic(s, str(media / spec.image), cfg.image_box)
            elif spec.code:
                box = cfg.code_box if spec.bullets else cfg.code_box_full
                _add_code(s, spec.code, spec.code_caption, theme, box, cfg.code_style)
        _add_materials(s, spec.materials, theme)
        if spec.notes:
            _set_notes(s, spec.notes)

    out = cfg.out_dir / f"{cfg.out_name}.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    _provenance(prs, cfg)
    prs.save(str(out))
    method = hardlink_or_copy(out, cfg.downloads_link) if cfg.downloads_link else None
    log.success(
        f"pptx → {out}  ({len(content.slides)} slides)"
        + (f"  [{method} → {cfg.downloads_link}]" if method else "")
    )
    return out, method
