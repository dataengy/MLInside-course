#!/usr/bin/env python3
# Build a ~35-slide MLInside "Введение в dbt" deck on the MLInside template.
# Content = condensed/de-duplicated from the 84-slide source (per dedup-plan.md).
# Design = inherited from the template's masters/layouts (Corbel, accent #2419FF).
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

TPL = "/Users/nk.myg/Downloads/MLinside шаблон презентаций.pptx"
OUT = "/Users/nk.myg/Downloads/MLInside_Введение-в-dbt.pptx"

prs = Presentation(TPL)

# ── remove the 6 example slides, keep masters/layouts/theme ───────────────────
# Drop the relationship too, else the old slide parts stay in the package and new
# slides collide on partname (slide1.xml…) → duplicate zip entries.
from pptx.oxml.ns import qn
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    prs.part.drop_rel(sldId.get(qn('r:id')))
    sldIdLst.remove(sldId)

# ── layouts by name ───────────────────────────────────────────────────────────
LAYOUT = {l.name: l for m in prs.slide_masters for l in m.slide_layouts}
L_TITLE   = LAYOUT["Титульный слайд"]        # CENTER_TITLE(0) + SUBTITLE(1) — designer's title (visible)
L_SECTION = LAYOUT["Заголовок и объект"]     # TITLE(0)  (used as block/section title)
L_CONTENT = LAYOUT["Title and body"]         # TITLE(0) + BODY(1)
L_CLOSING = LAYOUT["Титульный слайд"]        # reuse title layout for a symmetric closing
L_TABLE   = LAYOUT["Только заголовок"]        # TITLE(0) only — for slides with a custom table

ACCENT = RGBColor(0x24, 0x19, 0xFF)   # MLInside indigo accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
ALT    = RGBColor(0xF2, 0xF2, 0xF7)   # light indigo-tinted alt row

def _cell(cell, text, *, bold, color, size, fill):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.text = text
    for run in p.runs:
        run.font.name = "Corbel"; run.font.size = Pt(size)
        run.font.bold = bold; run.font.color.rgb = color

def add_table_slide(title, headers, rows, col_ratios=None):
    s = prs.slides.add_slide(L_TABLE)
    s.shapes.title.text = title
    left, top, width = Inches(0.45), Inches(1.5), Inches(12.43)
    height = Inches(0.45 + 0.4 * len(rows))
    tbl = s.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    tbl.first_row = True
    tbl.horz_banding = False
    if col_ratios:
        tot = sum(col_ratios)
        for i, rr in enumerate(col_ratios):
            tbl.columns[i].width = Emu(int(int(width) * rr / tot))
    for c, h in enumerate(headers):
        _cell(tbl.cell(0, c), h, bold=True, color=WHITE, size=13, fill=ACCENT)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            _cell(tbl.cell(r, c), val, bold=False, color=DARK, size=12,
                  fill=(WHITE if r % 2 else ALT))
    return s

def _set_body(slide, bullets):
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(bullets):
        text, lvl = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = lvl

def add_title(title, subtitle):
    s = prs.slides.add_slide(L_TITLE)
    s.shapes.title.text = title
    s.placeholders[1].text = subtitle
    return s

def add_section(title):
    s = prs.slides.add_slide(L_SECTION)
    s.shapes.title.text = title
    return s

def add_content(title, bullets):
    s = prs.slides.add_slide(L_CONTENT)
    s.shapes.title.text = title
    _set_body(s, bullets)
    return s

def add_closing(title, subtitle):
    s = prs.slides.add_slide(L_CLOSING)
    s.shapes.title.text = title
    s.placeholders[1].text = subtitle
    return s

# ── deck ──────────────────────────────────────────────────────────────────────
add_title("Введение в dbt", "MLInside 2026  ·  Николай Крупий, Data Engineer")

add_content("План лекции", [
    "Введение, предпосылки и место dbt в data-стеке",
    "Экосистема и продукты dbt (Core, Cloud, Semantic Layer)",
    "dbt Core: архитектура, проект, материализации, Jinja, команды",
    "Тестирование, документация, пакеты",
    "Запуск в продакшене: оркестрация и CI/CD",
    "Дополнительные возможности, AI и тренды",
])

# ── Раздел 1 ──
add_section("Введение и экосистема")

add_content("Что такое dbt", [
    "dbt = Data Build Tool — трансформация данных SQL-моделями",
    "Реализует «T» в ELT: преобразования внутри DWH",
    "SQL + Jinja + YAML вместо императивного кода",
    "Из коробки: версионирование, тесты, документация, DAG зависимостей",
    "Инструмент Analytics Engineering — код вместо ручных выгрузок",
])

add_content("Предпосылки и история", [
    "Проблемы классического ETL: хрупкие пайплайны, дубли логики",
    "Появление облачных DWH: Snowflake, BigQuery, Redshift",
    "ELT вытеснил ETL — грузим сырое, трансформируем в DWH",
    "dbt (Fishtown Analytics, 2016) — SQL-first трансформации как код",
    "Пишется строчными: «dbt», а не «DBT»",
])

add_content("Место dbt и роль Analytics Engineer", [
    "Поток данных: ingestion → transform (dbt) → BI / ML",
    "dbt закрывает слой трансформаций (Analytics Engineering)",
    "Роли: Data Engineer · Analytics Engineer · Data Analyst",
    "Analytics Engineer соединяет инженерные практики и аналитику",
    "Код, ревью, тесты и документация приходят в аналитику",
])

add_content("Продукты экосистемы dbt", [
    "dbt Core — open-source CLI, основа всего",
    "dbt Cloud — управляемая платформа (IDE, оркестрация, доки)",
    "Semantic Layer / MetricFlow — единый слой метрик",
    "dbt packages — переиспользуемые макросы, тесты, модели",
    "dbt MCP — доступ AI/агентов к проекту и метрикам",
])

add_table_slide("dbt Core vs dbt Cloud",
    ["", "dbt Core", "dbt Cloud"],
    [
        ["Тип", "Open-source CLI", "SaaS-платформа"],
        ["Хостинг", "Self-hosted", "Управляемый (dbt Labs)"],
        ["IDE", "Своя (VS Code и др.)", "Встроенная веб-IDE"],
        ["Оркестрация", "Airflow / Dagster / cron", "Встроенный планировщик"],
        ["Возможности", "Пакеты, полный контроль", "CI, доки, роли и доступы"],
        ["Стоимость", "Бесплатно", "Подписка"],
    ],
    col_ratios=[2, 4, 4])

# ── Раздел 2 ──
add_section("dbt Core: архитектура и проект")

add_content("dbt Core: что это и архитектура", [
    "CLI-приложение: компилирует SQL-модели с Jinja",
    "Автоматически управляет зависимостями и строит DAG",
    "Выполняет модели в нужном порядке в целевой БД",
    "Code-based ELT: модели хранятся как код в git",
    "Артефакты компиляции — в каталоге target/",
])

add_content("dbt как код: структура проекта", [
    "models/ — SQL-модели (.sql) + метаданные (schema.yml)",
    "dbt_project.yml — главный конфиг проекта",
    "profiles.yml — подключение к БД (обычно ~/.dbt/)",
    "seeds/, snapshots/, macros/, tests/ — прочие сущности",
    "Всё версионируется в git, ревьюится через MR/PR",
])

add_content("SQL + компиляция (обёртка кода)", [
    "Пишем select-логику модели в models/model.sql",
    "ref() и source() задают зависимости вместо хардкода имён",
    "dbt компилирует Jinja → чистый SQL под конкретную БД",
    "И оборачивает в DDL: create view / create table as select",
    "Скомпилированный код — в target/compiled и target/run",
])

add_content("DAG зависимостей", [
    "DAG — направленный ациклический граф моделей",
    "Строится автоматически из ref() между моделями",
    "Определяет порядок запуска и распараллеливание",
    "Даёт lineage: откуда и куда идут данные",
    "Основа для инкрементальных и частичных запусков",
])

add_table_slide("Jinja и макросы",
    ["Конструкция", "Пример", "Назначение"],
    [
        ["{{ ... }}", "{{ ref('stg_orders') }}", "выражения и вызовы функций"],
        ["{% ... %}", "{% for c in cols %}", "управляющие конструкции (if/for/set)"],
        ["ref()", "ref('model_b')", "зависимости между моделями"],
        ["source()", "source('raw','orders')", "ссылка на внешние таблицы"],
        ["config()", "{{ config(materialized='table') }}", "настройки модели"],
        ["макросы", "{{ my_macro(col) }}", "переиспользуемые функции"],
    ],
    col_ratios=[2, 4, 4])

add_table_slide("Материализации",
    ["Тип", "Что делает", "Где хранится"],
    [
        ["view", "CREATE VIEW при каждом run", "вьюха в БД"],
        ["table", "пересоздаёт таблицу целиком", "таблица в БД"],
        ["incremental", "дозагружает новые/изменённые строки", "таблица в БД"],
        ["ephemeral", "встраивается как CTE в другие модели", "нет объекта в БД"],
    ],
    col_ratios=[2, 5, 3])

add_content("Сущности dbt", [
    "models — основная логика (SQL/Jinja), строит таблицы/вьюхи",
    "sources — декларация внешних таблиц + freshness-проверки",
    "seeds — небольшие CSV, загружаемые как таблицы",
    "snapshots — историзация изменений (SCD2)",
    "tests, macros, exposures — качество, переиспользование, витрины",
])

add_content("Конфиги: dbt_project.yml и profiles.yml", [
    "dbt_project.yml — имя проекта, пути, материализации по умолчанию",
    ("Настройки моделей по папкам, vars, hooks", 1),
    "profiles.yml — подключение к БД (target, host, schema, креды)",
    ("Обычно ~/.dbt/profiles.yml; секреты через env_var()", 1),
    "Разделение: логика проекта отдельно от параметров окружения",
])

add_table_slide("Основные команды dbt",
    ["Команда", "Что делает"],
    [
        ["dbt run", "компилирует и выполняет модели"],
        ["dbt build", "run + test + seed + snapshot в порядке DAG"],
        ["dbt test", "тесты данных и unit-тесты"],
        ["dbt seed", "загружает CSV из seeds/ как таблицы"],
        ["dbt snapshot", "историзация изменений (SCD2)"],
        ["dbt docs generate", "генерирует документацию и lineage"],
        ["--select / --exclude", "выбор моделей (state:modified и др.)"],
    ],
    col_ratios=[3, 7])

add_table_slide("Адаптеры",
    ["Платформа", "Пакет", "Статус"],
    [
        ["PostgreSQL", "dbt-postgres", "Trusted"],
        ["Snowflake", "dbt-snowflake", "Trusted"],
        ["BigQuery", "dbt-bigquery", "Trusted"],
        ["Redshift", "dbt-redshift", "Trusted"],
        ["Databricks", "dbt-databricks", "Trusted"],
        ["ClickHouse", "dbt-clickhouse", "Community"],
        ["DuckDB", "dbt-duckdb", "Community"],
    ],
    col_ratios=[4, 4, 2])

# ── Раздел 3 ──
add_section("Тестирование, документация, пакеты")

add_content("Документирование в dbt", [
    "Документация генерируется автоматически: dbt docs generate",
    "Описания моделей и колонок берутся из schema.yml",
    "Строится интерактивный сайт с поиском и DAG/lineage",
    "Docs всегда синхронны с кодом — обновляются при каждом run",
    "Полезно аналитикам и стейкхолдерам как каталог данных",
])

add_content("Тестирование", [
    "Generic-тесты: not_null, unique, accepted_values, relationships",
    "Задаются декларативно в schema.yml на колонках",
    "Singular-тесты — произвольный SQL, возвращающий «плохие» строки",
    "Unit-тесты — проверяют логику модели на фиктивных входах",
    "Пакеты расширяют набор: dbt_utils, dbt_expectations",
])

add_content("Пакеты dbt: подключение", [
    "Пакет = отдельный dbt-проект с макросами/моделями/тестами",
    "Подключение: packages.yml + dbt deps",
    "Каталог: hub.getdbt.com; также git и локальные пути",
    "Свой пакет — чтобы переиспользовать код между проектами",
    "Экономят время и стандартизируют практики",
])

add_table_slide("Ключевые пакеты",
    ["Пакет", "Назначение"],
    [
        ["dbt_utils", "базовые макросы и generic-тесты"],
        ["dbt_expectations", "data-quality тесты (Great Expectations)"],
        ["dbt-date", "унифицированная работа с датами"],
        ["audit_helper", "сравнение моделей при рефакторинге"],
        ["elementary", "мониторинг качества и аномалий"],
        ["project_evaluator", "аудит структуры проекта"],
        ["dbtVault / AutomateDV", "автоматизация Data Vault"],
    ],
    col_ratios=[3, 7])

# ── Раздел 4 ──
add_section("Запуск в продакшене")

add_content("Как запускать dbt в проде", [
    "Цель: надёжно и повторяемо выполнять build/run/test",
    "По расписанию или по событию, с логами и алертами",
    "Ретраи и идемпотентность запусков",
    "Изоляция окружений: dev / staging / prod (разные targets)",
    "Версионирование и деплой проекта через CI",
])

add_content("Оркестрация: Airflow и Dagster", [
    "Airflow — классический оркестратор; dbt как набор задач/оператор",
    "Cosmos разворачивает dbt-DAG в Airflow-DAG",
    "Dagster — нативная интеграция: модели dbt как assets",
    "Единый граф data-ассетов, наблюдаемость, партиции",
    "Backfill — параметризованные прогоны за прошлые окна",
])

add_content("CI/CD для dbt", [
    "pre-commit — ловит ошибки до пуша (форматирование, стиль)",
    "GitLab CI / GitHub Actions — автопроверка перед merge",
    "SQLFluff — линтер и автоформаттер SQL (учитывает Jinja)",
    "Slim CI — гоняем только изменённые узлы DAG (state:modified)",
    "Быстрее и дешевле, чем прогон всего проекта",
])

# ── Раздел 5 ──
add_section("Дополнительные возможности и тренды")

add_content("Python-модели, микробатчинг, каталоги", [
    "Python-модели — когда SQL неудобен (ML-фичи, сложная логика)",
    "Микробатчинг — инкрементальная обработка окнами времени",
    "Артефакты: manifest.json, catalog.json, run_results.json",
    "Интеграция с Data Catalogs (DataHub, OpenMetadata и др.)",
    "Единый проект как источник метаданных и lineage",
])

add_content("Semantic Layer / MetricFlow", [
    "Единое определение метрик — один раз, в коде",
    "MetricFlow генерирует корректный SQL по запросу метрики",
    "Убирает расхождения метрик между дашбордами",
    "Метрики доступны из BI, ноутбуков и через API",
    "Консистентность и переиспользование бизнес-логики",
])

add_content("dbt и AI", [
    "dbt Copilot — генерация SQL, тестов, документации, метрик",
    "dbt MCP — стандартный доступ LLM/агентов к моделям и метрикам",
    "context7 / MCP — подключение IDE и AI к контексту проекта",
    "AI ускоряет рутину, но код остаётся ревьюируемым",
])

add_content("Рынок и тренды", [
    "dbt Fusion Engine — новый быстрый движок (бета, 2025)",
    "dbt Labs + Fivetran — консолидация data-стека",
    "SQLMesh — альтернатива с виртуальными окружениями",
    "OpenDBT — open-source расширения и локальные материализации",
    "Экосистема быстро развивается — следим за релизами",
])

add_content("Полезные ссылки", [
    "docs.getdbt.com — официальная документация",
    "hub.getdbt.com — каталог пакетов",
    "courses.getdbt.com — dbt Learn (бесплатные курсы)",
    "discourse.getdbt.com — обсуждения и рецепты",
    "dbt Community Slack — помощь и нетворкинг",
])

add_closing("Спасибо за внимание!",
    "Николай Крупий · Data Engineer · hnkovr@gmail.com · t.me/NikolayKrupiy")

prs.save(OUT)
print("SAVED:", OUT)
print("slides:", len(prs.slides._sldIdLst))
