#!/usr/bin/env python3
# Build v2 of the MLInside "Введение в dbt" deck:
#   • content condensed from the 84-slide source (per dedup-plan.md), on the MLInside template
#   • real diagrams/screenshots reused from the source deck's media, on content slides
#   • speaker notes (речь докладчика) on EVERY slide
# Media is extracted from the source pptx at build time → reproducible & self-contained.
import os, zipfile, tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

TPL = "/Users/nk.myg/Downloads/MLinside шаблон презентаций.pptx"
SRC = "/Users/nk.myg/Library/Mobile Documents/com~apple~CloudDocs/_2025-11-ВШЭ_ВНЕШНЕЕ_ОБУЧЕНИЕ/_!_HSU_202612/NEW_ВШЭ_Семинар9-Практикум по dbt.pptx"
OUT = "/Users/nk.myg/Downloads/MLInside_Введение-в-dbt_v2.pptx"

# ── extract source media (reproducible) ───────────────────────────────────────
MEDIA = os.path.expandvars("$CLAUDE_JOB_DIR/tmp/media")
if not (MEDIA and os.path.isdir(MEDIA)):
    MEDIA = tempfile.mkdtemp(prefix="dbt_media_")
    with zipfile.ZipFile(SRC) as z:
        for n in z.namelist():
            if n.startswith("ppt/media/"):
                open(os.path.join(MEDIA, os.path.basename(n)), "wb").write(z.read(n))
def img(name): return os.path.join(MEDIA, name)

# ── base deck from template ───────────────────────────────────────────────────
prs = Presentation(TPL)
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    prs.part.drop_rel(sldId.get(qn('r:id'))); sldIdLst.remove(sldId)
LAYOUT = {l.name: l for m in prs.slide_masters for l in m.slide_layouts}
L_TITLE   = LAYOUT["Титульный слайд"]
L_SECTION = LAYOUT["Заголовок и объект"]
L_CONTENT = LAYOUT["Title and body"]
L_TABLE   = LAYOUT["Только заголовок"]

ACCENT = RGBColor(0x24, 0x19, 0xFF); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A); ALT   = RGBColor(0xF2, 0xF2, 0xF7)

def _notes(s, text):
    if text:
        s.notes_slide.notes_text_frame.text = text

def _set_body(slide, bullets):
    tf = slide.placeholders[1].text_frame; tf.word_wrap = True; tf.clear()
    for i, item in enumerate(bullets):
        text, lvl = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.level = lvl

def _add_pic(slide, path, box=(7.35, 1.65, 5.5, 5.0)):
    bl, bt, bw, bh = [Inches(v) for v in box]
    iw, ih = Image.open(path).size
    if iw / ih > bw / bh:
        w = bw; h = int(bw * ih / iw)
    else:
        h = bh; w = int(bh * iw / ih)
    slide.shapes.add_picture(path, bl + (bw - w) // 2, bt + (bh - h) // 2, width=w, height=h)

def add_title(title, subtitle, notes=None):
    s = prs.slides.add_slide(L_TITLE)
    s.shapes.title.text = title; s.placeholders[1].text = subtitle
    _notes(s, notes); return s

def add_section(title, notes=None):
    s = prs.slides.add_slide(L_SECTION); s.shapes.title.text = title
    _notes(s, notes); return s

def add_content(title, bullets, image=None, notes=None):
    s = prs.slides.add_slide(L_CONTENT); s.shapes.title.text = title
    if image:
        ph = s.placeholders[1]
        # Materialize INHERITED geometry first — setting only .width creates a partial
        # xfrm (offset defaults to 0,0) and the body jumps under the title. Capture all,
        # then assign together so the body keeps its place and only narrows.
        l0, t0, h0 = ph.left, ph.top, ph.height
        ph.left, ph.top, ph.height, ph.width = l0, t0, h0, Inches(6.7)
    _set_body(s, bullets)
    if image:
        _add_pic(s, img(image))
    _notes(s, notes); return s

def _cell(cell, text, *, bold, color, size, fill):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.12); tf.margin_top = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.text = text
    for r in p.runs:
        r.font.name = "Corbel"; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

def add_table_slide(title, headers, rows, col_ratios=None, notes=None):
    s = prs.slides.add_slide(L_TABLE); s.shapes.title.text = title
    left, top, width = Inches(0.45), Inches(1.5), Inches(12.43)
    tbl = s.shapes.add_table(len(rows) + 1, len(headers), left, top, width,
                             Inches(0.45 + 0.4 * len(rows))).table
    tbl.first_row = True; tbl.horz_banding = False
    if col_ratios:
        tot = sum(col_ratios)
        for i, rr in enumerate(col_ratios):
            tbl.columns[i].width = Emu(int(int(width) * rr / tot))
    for c, h in enumerate(headers):
        _cell(tbl.cell(0, c), h, bold=True, color=WHITE, size=13, fill=ACCENT)
    for r, row in enumerate(rows, 1):
        for c, v in enumerate(row):
            _cell(tbl.cell(r, c), v, bold=False, color=DARK, size=12, fill=(WHITE if r % 2 else ALT))
    _notes(s, notes); return s

def add_closing(title, subtitle, notes=None):
    s = prs.slides.add_slide(L_TITLE); s.shapes.title.text = title
    s.placeholders[1].text = subtitle; _notes(s, notes); return s

# ══ DECK ══════════════════════════════════════════════════════════════════════
add_title("Введение в dbt", "MLInside 2026  ·  Николай Крупий, Data Engineer",
    notes="Приветствие. Сегодня разбираем dbt — инструмент трансформации данных, который стал стандартом Analytics Engineering. Цель лекции — понять, что такое dbt, как устроен dbt Core, и как его применяют в реальных проектах: от моделей и тестов до запуска в продакшене.")

add_content("План лекции", [
    "Введение, предпосылки и место dbt в data-стеке",
    "Экосистема и продукты dbt (Core, Cloud, Semantic Layer)",
    "dbt Core: архитектура, проект, материализации, Jinja, команды",
    "Тестирование, документация, пакеты",
    "Запуск в продакшене: оркестрация и CI/CD",
    "Дополнительные возможности, AI и тренды",
], notes="Коротко пройдёмся по плану. Начнём с того, зачем вообще dbt появился, затем нырнём в устройство dbt Core — это ядро, потом посмотрим тестирование и пакеты, и закончим продакшеном и трендами. Вопросы удобнее задавать по ходу.")

add_section("Введение и экосистема",
    notes="Первый блок — вводный. Разберёмся с базовыми понятиями: что такое dbt, откуда он взялся, какое место занимает в data-платформе и какие продукты есть в экосистеме.")

add_content("Что такое dbt", [
    "dbt = Data Build Tool — трансформация данных SQL-моделями",
    "Реализует «T» в ELT: преобразования внутри DWH",
    "SQL + Jinja + YAML вместо императивного кода",
    "Из коробки: версионирование, тесты, документация, DAG зависимостей",
    "Инструмент Analytics Engineering — код вместо ручных выгрузок",
], image="image28.png",
    notes="dbt отвечает за букву T в ELT — трансформацию. Данные уже загружены в хранилище, а dbt превращает сырые таблицы в чистые витрины с помощью обычного SQL. Ключевая идея: аналитическая логика становится кодом — с git, ревью, тестами и автодокументацией. На схеме — типичный конвейр: источники, загрузка, трансформация в dbt, витрины.")

add_content("Предпосылки и история", [
    "Проблемы классического ETL: хрупкие пайплайны, дубли логики",
    "Появление облачных DWH: Snowflake, BigQuery, Redshift",
    "ELT вытеснил ETL — грузим сырое, трансформируем в DWH",
    "dbt (Fishtown Analytics, 2016) — SQL-first трансформации как код",
    "Пишется строчными: «dbt», а не «DBT»",
], image="image23.jpeg",
    notes="Раньше трансформации жили в тяжёлых ETL-инструментах и питоновских скриптах. С приходом дешёвых мощных облачных хранилищ логичнее стало грузить сырые данные как есть и трансформировать уже внутри БД — это подход ELT. dbt появился в 2016 году и закрыл именно слой T. И маленькая деталь: название пишется строчными буквами.")

add_content("Место dbt и роль Analytics Engineer", [
    "Поток данных: ingestion → transform (dbt) → BI / ML",
    "dbt закрывает слой трансформаций (Analytics Engineering)",
    "Роли: Data Engineer · Analytics Engineer · Data Analyst",
    "Analytics Engineer соединяет инженерные практики и аналитику",
    "Код, ревью, тесты и документация приходят в аналитику",
], image="image24.png",
    notes="dbt породил отдельную роль — Analytics Engineer. Data Engineer отвечает за инфраструктуру и загрузку, Data Analyst — за инсайты и дашборды, а Analytics Engineer посередине: строит надёжные, протестированные модели данных, применяя инженерные практики к аналитике. На схеме — как делятся зоны ответственности между тремя ролями.")

add_content("Продукты экосистемы dbt", [
    "dbt Core — open-source CLI, основа всего",
    "dbt Cloud — управляемая платформа (IDE, оркестрация, доки)",
    "Semantic Layer / MetricFlow — единый слой метрик",
    "dbt packages — переиспользуемые макросы, тесты, модели",
    "dbt MCP — доступ AI/агентов к проекту и метрикам",
], image="image27.png",
    notes="Экосистема шире, чем один CLI. В центре — dbt Core, бесплатное ядро. Вокруг: dbt Cloud как управляемая платформа, Semantic Layer для метрик, пакеты для переиспользования кода и свежий dbt MCP для интеграции с ИИ. Сегодня фокус на Core, но важно видеть всю картину.")

add_table_slide("dbt Core vs dbt Cloud",
    ["", "dbt Core", "dbt Cloud"],
    [["Тип", "Open-source CLI", "SaaS-платформа"],
     ["Хостинг", "Self-hosted", "Управляемый (dbt Labs)"],
     ["IDE", "Своя (VS Code и др.)", "Встроенная веб-IDE"],
     ["Оркестрация", "Airflow / Dagster / cron", "Встроенный планировщик"],
     ["Возможности", "Пакеты, полный контроль", "CI, доки, роли и доступы"],
     ["Стоимость", "Бесплатно", "Подписка"]],
    col_ratios=[2, 4, 4],
    notes="Главное различие: Core — это бесплатный инструмент командной строки, который вы разворачиваете и оркеструете сами; Cloud — платная управляемая среда со всем готовым: веб-IDE, планировщик, CI, хостинг документации. Для обучения и полного контроля берём Core, для быстрого старта без инфраструктуры — Cloud.")

add_section("dbt Core: архитектура и проект",
    notes="Переходим к ядру. Разберём, как устроен dbt Core изнутри: компиляция, DAG, структура проекта, материализации, Jinja и основные команды. Это самый практический блок.")

add_content("dbt Core: что это и архитектура", [
    "CLI-приложение: компилирует SQL-модели с Jinja",
    "Автоматически управляет зависимостями и строит DAG",
    "Выполняет модели в нужном порядке в целевой БД",
    "Code-based ELT: модели хранятся как код в git",
    "Артефакты компиляции — в каталоге target/",
], image="image33.png",
    notes="dbt Core — это приложение командной строки. Вы пишете SQL-модели с вставками Jinja, а dbt компилирует их в чистый SQL, строит граф зависимостей и выполняет модели в правильном порядке в вашей БД. На экране — типичный вывод dbt run: видно, как модели выполняются одна за другой.")

add_content("dbt как код: структура проекта", [
    "models/ — SQL-модели (.sql) + метаданные (schema.yml)",
    "dbt_project.yml — главный конфиг проекта",
    "profiles.yml — подключение к БД (обычно ~/.dbt/)",
    "seeds/, snapshots/, macros/, tests/ — прочие сущности",
    "Всё версионируется в git, ревьюится через MR/PR",
], image="image30.png",
    notes="Проект dbt — это структурированная папка в git. В models лежат SQL-файлы и рядом yml с описаниями и тестами. Есть главный конфиг проекта и отдельный файл подключения к БД. Поскольку всё это код, работают привычные практики: ветки, ревью, история. На скриншоте — проект dbt в IDE.")

add_content("SQL + компиляция (обёртка кода)", [
    "Пишем select-логику модели в models/model.sql",
    "ref() и source() задают зависимости вместо хардкода имён",
    "dbt компилирует Jinja → чистый SQL под конкретную БД",
    "И оборачивает в DDL: create view / create table as select",
    "Скомпилированный код — в target/compiled и target/run",
], image="image40.png",
    notes="Вы пишете только select — саму логику. Ссылки на другие модели даёте через ref, а не хардкодите имена таблиц — именно это позволяет dbt построить граф. При запуске dbt раскрывает Jinja в чистый SQL и оборачивает его в create view или create table. Готовый SQL всегда можно посмотреть в папке target.")

add_content("DAG зависимостей", [
    "DAG — направленный ациклический граф моделей",
    "Строится автоматически из ref() между моделями",
    "Определяет порядок запуска и распараллеливание",
    "Даёт lineage: откуда и куда идут данные",
    "Основа для инкрементальных и частичных запусков",
], image="image37.png",
    notes="Из вызовов ref dbt автоматически строит DAG — граф зависимостей. Он определяет порядок выполнения и что можно запускать параллельно. Бонусом получаем lineage — наглядную схему, откуда и куда текут данные. На картинке — простой граф: staging-модели питают промежуточные, а те — витрины.")

add_table_slide("Jinja и макросы",
    ["Конструкция", "Пример", "Назначение"],
    [["{{ ... }}", "{{ ref('stg_orders') }}", "выражения и вызовы функций"],
     ["{% ... %}", "{% for c in cols %}", "управляющие конструкции (if/for/set)"],
     ["ref()", "ref('model_b')", "зависимости между моделями"],
     ["source()", "source('raw','orders')", "ссылка на внешние таблицы"],
     ["config()", "{{ config(materialized='table') }}", "настройки модели"],
     ["макросы", "{{ my_macro(col) }}", "переиспользуемые функции"]],
    col_ratios=[2, 4, 4],
    notes="Jinja — это шаблонизатор, который делает SQL динамическим. Двойные фигурные скобки — это выражения и вызовы функций, процентные — управляющие конструкции вроде циклов и условий. Ключевые функции — ref, source, config. А макросы — это ваши собственные переиспользуемые функции, как в обычном программировании.")

add_table_slide("Материализации",
    ["Тип", "Что делает", "Где хранится"],
    [["view", "CREATE VIEW при каждом run", "вьюха в БД"],
     ["table", "пересоздаёт таблицу целиком", "таблица в БД"],
     ["incremental", "дозагружает новые/изменённые строки", "таблица в БД"],
     ["ephemeral", "встраивается как CTE в другие модели", "нет объекта в БД"]],
    col_ratios=[2, 5, 3],
    notes="Материализация — это то, во что превращается ваша модель в БД. view — просто представление, всегда свежее, но считается на лету. table — физическая таблица, быстрые чтения ценой полного пересчёта. incremental — дозагружает только новое, ключ к большим данным. ephemeral — не создаёт объект, а встраивается в другие модели как CTE.")

add_content("Сущности dbt", [
    "models — основная логика (SQL/Jinja), строит таблицы/вьюхи",
    "sources — декларация внешних таблиц + freshness-проверки",
    "seeds — небольшие CSV, загружаемые как таблицы",
    "snapshots — историзация изменений (SCD2)",
    "tests, macros, exposures — качество, переиспользование, витрины",
], image="image41.png",
    notes="Кроме моделей в dbt есть и другие сущности. sources описывают внешние таблицы и умеют проверять свежесть данных. seeds — маленькие справочники из CSV. snapshots историзируют изменения по типу SCD2. Плюс тесты, макросы и exposures. Всё это — части одного проекта.")

add_content("Конфиги: dbt_project.yml и profiles.yml", [
    "dbt_project.yml — имя проекта, пути, материализации по умолчанию",
    ("Настройки моделей по папкам, vars, hooks", 1),
    "profiles.yml — подключение к БД (target, host, schema, креды)",
    ("Обычно ~/.dbt/profiles.yml; секреты через env_var()", 1),
    "Разделение: логика проекта отдельно от параметров окружения",
], image="image42.png",
    notes="Два ключевых конфига. dbt_project.yml описывает сам проект: имя, пути, материализации по умолчанию, переменные. profiles.yml — это подключение к БД, он живёт отдельно, обычно в домашней папке, а секреты подтягиваются через env_var. Такое разделение позволяет один и тот же проект гонять в разных окружениях.")

add_table_slide("Основные команды dbt",
    ["Команда", "Что делает"],
    [["dbt run", "компилирует и выполняет модели"],
     ["dbt build", "run + test + seed + snapshot в порядке DAG"],
     ["dbt test", "тесты данных и unit-тесты"],
     ["dbt seed", "загружает CSV из seeds/ как таблицы"],
     ["dbt snapshot", "историзация изменений (SCD2)"],
     ["dbt docs generate", "генерирует документацию и lineage"],
     ["--select / --exclude", "выбор моделей (state:modified и др.)"]],
    col_ratios=[3, 7],
    notes="Команд немного. dbt run выполняет модели, dbt test гоняет тесты, а dbt build делает всё сразу в правильном порядке по DAG — это рабочая лошадка. Отдельно — seed, snapshot и генерация документации. И почти к любой команде можно добавить select, чтобы запустить не весь проект, а только нужные модели.")

add_table_slide("Адаптеры",
    ["Платформа", "Пакет", "Статус"],
    [["PostgreSQL", "dbt-postgres", "Trusted"], ["Snowflake", "dbt-snowflake", "Trusted"],
     ["BigQuery", "dbt-bigquery", "Trusted"], ["Redshift", "dbt-redshift", "Trusted"],
     ["Databricks", "dbt-databricks", "Trusted"], ["ClickHouse", "dbt-clickhouse", "Community"],
     ["DuckDB", "dbt-duckdb", "Community"]],
    col_ratios=[4, 4, 2],
    notes="dbt работает с БД через адаптеры. Для крупных облачных хранилищ есть официальные, доверенные адаптеры. Для многих других, включая ClickHouse и DuckDB, — качественные community-адаптеры. Ставится нужный адаптер отдельным пакетом, и один и тот же проект в разумных пределах переносим между платформами.")

add_section("Тестирование, документация, пакеты",
    notes="Следующий блок — про качество и переиспользование. Как тестировать данные, как получать документацию бесплатно и как подключать готовые пакеты, чтобы не изобретать велосипед.")

add_content("Документирование в dbt", [
    "Документация генерируется автоматически: dbt docs generate",
    "Описания моделей и колонок берутся из schema.yml",
    "Строится интерактивный сайт с поиском и DAG/lineage",
    "Docs всегда синхронны с кодом — обновляются при каждом run",
    "Полезно аналитикам и стейкхолдерам как каталог данных",
], image="image45.png",
    notes="Документация в dbt — почти бесплатный бонус. Вы пишете описания в yml рядом с моделями, а dbt docs generate собирает интерактивный сайт с поиском и графом lineage. Главное — он всегда актуален, потому что генерируется из того же кода. Для аналитиков это фактически каталог данных.")

add_content("Тестирование", [
    "Generic-тесты: not_null, unique, accepted_values, relationships",
    "Задаются декларативно в schema.yml на колонках",
    "Singular-тесты — произвольный SQL, возвращающий «плохие» строки",
    "Unit-тесты — проверяют логику модели на фиктивных входах",
    "Пакеты расширяют набор: dbt_utils, dbt_expectations",
], image="image48.png",
    notes="Тесты в dbt бывают двух видов. Generic — готовые проверки вроде not_null и unique, которые вы просто вешаете на колонку в yml. Singular — произвольный SQL, который должен вернуть ноль строк. Плюс относительно новые unit-тесты для логики моделей. А пакеты вроде dbt_expectations сильно расширяют арсенал проверок.")

add_content("Пакеты dbt: подключение", [
    "Пакет = отдельный dbt-проект с макросами/моделями/тестами",
    "Подключение: packages.yml + dbt deps",
    "Каталог: hub.getdbt.com; также git и локальные пути",
    "Свой пакет — чтобы переиспользовать код между проектами",
    "Экономят время и стандартизируют практики",
], image="image46.jpg",
    notes="Пакет — это просто другой dbt-проект, который вы подключаете. Прописываете его в packages.yml, выполняете dbt deps — и получаете чужие макросы, тесты и модели. Основной каталог — hub.getdbt.com, вот он на экране. Можно и свой пакет собрать, чтобы переиспользовать код между проектами команды.")

add_table_slide("Ключевые пакеты",
    ["Пакет", "Назначение"],
    [["dbt_utils", "базовые макросы и generic-тесты"],
     ["dbt_expectations", "data-quality тесты (Great Expectations)"],
     ["dbt-date", "унифицированная работа с датами"],
     ["audit_helper", "сравнение моделей при рефакторинге"],
     ["elementary", "мониторинг качества и аномалий"],
     ["project_evaluator", "аудит структуры проекта"],
     ["dbtVault / AutomateDV", "автоматизация Data Vault"]],
    col_ratios=[3, 7],
    notes="Несколько пакетов, которые стоит знать. dbt_utils — практически стандарт, базовые макросы и тесты. dbt_expectations добавляет проверки качества данных. audit_helper незаменим при рефакторинге. elementary следит за аномалиями. project_evaluator проверяет, что проект соответствует лучшим практикам. С них удобно начинать.")

add_section("Запуск в продакшене",
    notes="Теперь — как всё это запускать не с ноутбука, а надёжно в продакшене: оркестрация, расписания, окружения и CI/CD.")

add_content("Как запускать dbt в проде", [
    "Цель: надёжно и повторяемо выполнять build/run/test",
    "По расписанию или по событию, с логами и алертами",
    "Ретраи и идемпотентность запусков",
    "Изоляция окружений: dev / staging / prod (разные targets)",
    "Версионирование и деплой проекта через CI",
], image="image43.png",
    notes="В продакшене нам нужна надёжность и повторяемость: запуск по расписанию или событию, логи, алерты, ретраи. Обязательно разделяем окружения — dev, staging, prod — через разные targets в профиле. И сам проект деплоится через CI, а не копируется руками.")

add_content("Оркестрация: Airflow и Dagster", [
    "Airflow — классический оркестратор; dbt как набор задач/оператор",
    "Cosmos разворачивает dbt-DAG в Airflow-DAG",
    "Dagster — нативная интеграция: модели dbt как assets",
    "Единый граф data-ассетов, наблюдаемость, партиции",
    "Backfill — параметризованные прогоны за прошлые окна",
], image="image58.png",
    notes="Два популярных оркестратора. Airflow — классика: dbt запускается как задачи, а библиотека Cosmos умеет разворачивать dbt-граф прямо в Airflow-DAG. Dagster — более современный подход: модели dbt становятся ассетами в едином графе с хорошей наблюдаемостью и партициями. На экране — интерфейс Dagster с моделями dbt.")

add_content("CI/CD для dbt", [
    "pre-commit — ловит ошибки до пуша (форматирование, стиль)",
    "GitLab CI / GitHub Actions — автопроверка перед merge",
    "SQLFluff — линтер и автоформаттер SQL (учитывает Jinja)",
    "Slim CI — гоняем только изменённые узлы DAG (state:modified)",
    "Быстрее и дешевле, чем прогон всего проекта",
], image="image48.png",
    notes="CI/CD для dbt устроен как для обычного кода. pre-commit ловит ошибки ещё до пуша, CI проверяет изменения перед merge, SQLFluff следит за стилем SQL. Отдельно стоит выделить Slim CI: он прогоняет только изменённые модели и их зависимых, а не весь проект — это сильно экономит время и деньги.")

add_section("Дополнительные возможности и тренды",
    notes="Финальный блок — что ещё умеет dbt за пределами базового SQL, и куда движется экосистема: Python-модели, Semantic Layer, интеграция с ИИ и свежие тренды рынка.")

add_content("Python-модели, микробатчинг, каталоги", [
    "Python-модели — когда SQL неудобен (ML-фичи, сложная логика)",
    "Микробатчинг — инкрементальная обработка окнами времени",
    "Артефакты: manifest.json, catalog.json, run_results.json",
    "Интеграция с Data Catalogs (DataHub, OpenMetadata и др.)",
    "Единый проект как источник метаданных и lineage",
], image="image60.png",
    notes="Несколько продвинутых возможностей. Там, где SQL неудобен — например, ML-фичи — можно писать Python-модели. Микробатчинг помогает с большими инкрементальными нагрузками. А артефакты, которые dbt генерирует — manifest и catalog — становятся источником метаданных для внешних каталогов данных вроде DataHub.")

add_content("Semantic Layer / MetricFlow", [
    "Единое определение метрик — один раз, в коде",
    "MetricFlow генерирует корректный SQL по запросу метрики",
    "Убирает расхождения метрик между дашбордами",
    "Метрики доступны из BI, ноутбуков и через API",
    "Консистентность и переиспользование бизнес-логики",
], image="image62.png",
    notes="Больная тема любой аналитики — когда «выручка» в двух дашбордах считается по-разному. Semantic Layer решает это: метрика определяется один раз в коде, а движок MetricFlow сам генерирует корректный SQL под каждый запрос. Метрики доступны из BI, ноутбуков и по API — и везде считаются одинаково.")

add_content("dbt и AI", [
    "dbt Copilot — генерация SQL, тестов, документации, метрик",
    "dbt MCP — стандартный доступ LLM/агентов к моделям и метрикам",
    "context7 / MCP — подключение IDE и AI к контексту проекта",
    "AI ускоряет рутину, но код остаётся ревьюируемым",
], image="image64.png",
    notes="dbt активно интегрируется с ИИ. Copilot генерирует SQL, тесты и документацию. dbt MCP даёт языковым моделям и агентам стандартный, управляемый доступ к вашим моделям и метрикам. Важно: ИИ ускоряет рутину, но результат — это по-прежнему код, который проходит ревью и тесты.")

add_content("Рынок и тренды", [
    "dbt Fusion Engine — новый быстрый движок (бета, 2025)",
    "dbt Labs + Fivetran — консолидация data-стека",
    "SQLMesh — альтернатива с виртуальными окружениями",
    "OpenDBT — open-source расширения и локальные материализации",
    "Экосистема быстро развивается — следим за релизами",
], image="image67.png",
    notes="Пара слов о трендах. dbt Labs выпустили новый движок Fusion и объединяются с Fivetran — идёт консолидация стека. Появляются альтернативы вроде SQLMesh с интересной моделью виртуальных окружений. Экосистема живая и быстро меняется, поэтому полезно следить за релизами.")

add_content("Полезные ссылки", [
    "docs.getdbt.com — официальная документация",
    "hub.getdbt.com — каталог пакетов",
    "courses.getdbt.com — dbt Learn (бесплатные курсы)",
    "discourse.getdbt.com — обсуждения и рецепты",
    "dbt Community Slack — помощь и нетворкинг",
], notes="Куда идти дальше. Официальная документация очень хорошая, обязательно бесплатные курсы dbt Learn — они компактные и практичные. За готовыми рецептами — на Discourse и в Community Slack. Ссылки будут в материалах.")

add_closing("Спасибо за внимание!",
    "Николай Крупий · Data Engineer · hnkovr@gmail.com · t.me/NikolayKrupiy",
    notes="Спасибо! Кратко резюме: dbt делает трансформации кодом — с тестами, документацией и графом зависимостей из коробки, а вокруг ядра выросла целая экосистема. Готов ответить на вопросы, контакты на экране.")

prs.save(OUT)
print("SAVED:", OUT, "| slides:", len(prs.slides._sldIdLst))
